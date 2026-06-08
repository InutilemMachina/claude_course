"""Mini-env: native PPTX equation (OMML) beillesztés teszt.
LaTeX -> MathML (latex2mathml) -> OMML (MML2OMML.XSL) -> a:txBody-ba ágyazva.
Block: m:oMathPara saját bekezdésként; Inline: m:oMath a:r run-ok között.
Mindkettő mc:AlternateContent/a14:m wrapperben (DrawingML 2010 math).
"""
from pathlib import Path
import latex2mathml.converter as l2m
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

XSL = etree.XSLT(etree.parse(r"C:\Program Files\Microsoft Office\root\Office16\MML2OMML.XSL"))
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "m": "http://schemas.openxmlformats.org/officeDocument/2006/math",
    "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
    "a14": "http://schemas.microsoft.com/office/drawing/2010/main",
}


def tex_to_omath(tex: str):
    """LaTeX -> <m:oMath> lxml elem."""
    mml = l2m.convert(tex)
    omml = XSL(etree.fromstring(mml.encode())).getroot()   # <m:oMath>
    return omml


def _amc(inner_xml: str) -> str:
    """mc:AlternateContent/a14:m wrapper (Fallback = sima szöveg)."""
    return (
        f'<mc:AlternateContent xmlns:mc="{NS["mc"]}">'
        f'<mc:Choice xmlns:a14="{NS["a14"]}" Requires="a14"><a14:m>'
        f'{inner_xml}'
        f'</a14:m></mc:Choice>'
        f'<mc:Fallback><a:r xmlns:a="{NS["a"]}"><a:t>[képlet]</a:t></a:r></mc:Fallback>'
        f'</mc:AlternateContent>'
    )


def omath_xml(omath) -> str:
    return etree.tostring(omath, encoding="unicode")


def block_para(tex: str) -> str:
    """Block egyenlet saját bekezdésként (m:oMathPara)."""
    om = omath_xml(tex_to_omath(tex))
    para = f'<m:oMathPara xmlns:m="{NS["m"]}">{om}</m:oMathPara>'
    return (f'<a:p xmlns:a="{NS["a"]}"><a:pPr algn="ctr"/>{_amc(para)}</a:p>')


def inline_para(prefix: str, tex: str, suffix: str) -> str:
    """Szövegközi egyenlet: run + m:oMath + run egy bekezdésben."""
    om = omath_xml(tex_to_omath(tex))
    return (
        f'<a:p xmlns:a="{NS["a"]}">'
        f'<a:r><a:t>{prefix}</a:t></a:r>'
        f'{_amc(om)}'
        f'<a:r><a:t>{suffix}</a:t></a:r>'
        f'</a:p>'
    )


def add_paragraph_xml(text_frame, p_xml: str):
    """Nyers <a:p> XML hozzáfűzése a txBody-hoz."""
    txBody = text_frame._txBody
    p = etree.fromstring(p_xml)
    txBody.append(p)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Natív egyenlet-teszt (OMML):"
    tf.paragraphs[0].runs[0].font.size = Pt(20)

    # Inline: szöveg + $...$ + szöveg
    add_paragraph_xml(tf, inline_para(
        "ahol ", r"P_o", " és "))
    add_paragraph_xml(tf, inline_para(
        "a teljes nyomásviszony ", r"B = \frac{U}{2a}\sqrt{\frac{V_p}{A_c L_c}}", " a B-paraméter."))
    # Block
    add_paragraph_xml(tf, block_para(
        r"SM = \frac{(P_o/P_i)_{Surge} - (P_o/P_i)_{SA}}{(P_o/P_i)_{SA}}"))
    add_paragraph_xml(tf, block_para(r"\frac{\partial PR}{\partial m} \leq 0"))

    out = Path(__file__).parent / "test1.pptx"
    prs.save(str(out))
    print("Mentve:", out)


if __name__ == "__main__":
    main()
