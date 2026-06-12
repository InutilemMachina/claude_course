"""
02_mineru_to_catalog.py — MinerU-first ábra-kinyerő + katalógus-builder (v4 séma).

MinerU-t használja az összes PDF feldolgozásához, majd a _content_list.json-ból
automatikusan kitölti a figure_catalog.json v4 séma mezőit:
  - id, page, path, needs_crop        ← script (determinisztikus)
  - caption                            ← MinerU image_caption
  - text_context                       ← MinerU ±3 szomszédos szöveg-entry
  - keywords (draft)                   ← caption + section heading + text top-words
  - visual_content                     ← null (02b_figure_enricher skill tölti)

PPTX-hez python-pptx fallback (slides-on caption ritka, text_context a dia szövege).

Ez a kanonikus 02 forrás-feldolgozó. A korábbi PyMuPDF-fallback
(02_image_extraction.py) megszűnt (P2.3, MinerU-only). Előfeltétel: conda `mineru` env.

Usage:
    python scripts/02_mineru_to_catalog.py --week-dir test_outputs/atg/1_het
    python scripts/02_mineru_to_catalog.py --week-dir test_outputs/atg/1_het --source X.pdf
    python scripts/02_mineru_to_catalog.py --week-dir test_outputs/atg/1_het --dry-run
    python scripts/02_mineru_to_catalog.py --week-dir test_outputs/atg/1_het --workers 3
    python scripts/02_mineru_to_catalog.py --week-dir test_outputs/atg/1_het --backend pipeline --device cuda
    python scripts/02_mineru_to_catalog.py --week-dir test_outputs/atg/1_het --no-resume

CLI argumentumok:
  --workers N           Párhuzamos worker-ek száma (default: cpu_count // 2, min 1).
                        Minden forrás külön MinerU folyamatban fut (ThreadPoolExecutor).
  --backend BACKEND     MinerU -b flag: pipeline (CPU default), vlm-auto-engine (GPU), stb.
  --device DEVICE       MinerU -d flag, csak pipeline backendhez: cpu, cuda, cuda:0, stb.
  --vram-per-worker N   GPU VRAM GB/worker (csak pipeline). Default: auto (nvidia-smi / workers).
  --no-resume           Mindig újrafuttatja MinerU-t, még ha content_list.json már létezik is.
                        Default (resume=True): megszakított futás folytatása MinerU újrafuttatása nélkül.
"""

from __future__ import annotations

import argparse
import concurrent.futures
import json
import os
import shutil
import subprocess
import sys
import threading
from collections import Counter
from datetime import date
from pathlib import Path

try:
    from _encoding_fix import fix_stdout as _fix_stdout
    _fix_stdout()
except ImportError:
    pass

# ── Konfiguráció ──────────────────────────────────────────────────────────────
MINERU_ENV     = "mineru"
CONTEXT_WINDOW = 3       # ±N szomszédos szöveg-entry a text_context-hez
MIN_AREA       = 10_000  # px² alatt (PPTX-ben): kihagyva
SCHEMA_VERSION = 4

LANG_HINT = {
    "nagyi":  "latin",
    "magyar": "latin",
    "hu":     "latin",
}

STOPWORDS = {
    "a", "an", "the", "and", "or", "of", "in", "on", "at", "to", "is", "are",
    "was", "were", "be", "been", "being", "by", "for", "with", "as", "it", "its",
    "this", "that", "from", "not", "but", "if", "has", "have", "had", "do",
    "does", "did", "will", "would", "can", "could", "should", "may", "might",
    "which", "who", "when", "where", "how", "what", "all", "also", "both",
    "each", "more", "such", "than", "so", "up", "out", "into", "about",
    "after", "between", "through", "during", "before", "above", "below",
    "az", "és", "is", "van", "meg", "el", "be", "ki", "fel", "le",
    "ez", "egy", "de", "ha", "nem", "hogy", "mint", "vagy", "sem", "már",
}

# ── v4 séma (figure_catalog.json) ─────────────────────────────────────────────
ENTRY_DEFAULTS = {
    "id":               None,
    "page":             None,
    "path":             None,
    "needs_crop":       False,
    "caption":          None,
    "caption_verified": False,
    "visual_content":   None,
    "text_context":     None,
    "keywords":         None,
    "_status":          "un-processed",
    "notes":            [],
}

CATALOG_GUIDE_TEMPLATE = """\
# figure_catalog.json — Útmutató

**Részletes szabályok:** `.claude/skills/02b_figure_enricher.md`
Forrás: `scripts/02_mineru_to_catalog.py` (MinerU-first pipeline)
"""


# ── Catalog helper-ek ─────────────────────────────────────────────────────────

def load_citations(week_dir: Path) -> dict:
    cit_path = week_dir / "1_raw_inputs" / "citations.json"
    if not cit_path.exists():
        return {}
    data = json.loads(cit_path.read_text(encoding="utf-8"))
    return {v["filename"]: k for k, v in data.items()
            if k != "_meta" and v.get("filename")}


def new_catalog() -> dict:
    return {
        "_meta": {
            "schema_version": SCHEMA_VERSION,
            "last_updated": date.today().isoformat(),
            "_guide": "CATALOG_GUIDE.md",
        },
        "sources": {},
    }


def load_catalog(path: Path) -> dict:
    if not path.exists():
        return new_catalog()
    raw = json.loads(path.read_text(encoding="utf-8"))
    schema = raw.get("_meta", {}).get("schema_version") if isinstance(raw, dict) else None
    if schema != SCHEMA_VERSION:
        sys.exit(f"HIBA: nem v{SCHEMA_VERSION} séma ({schema}) — wipe szükséges.")
    for src_data in raw.get("sources", {}).values():
        for e in src_data.get("figures", []):
            for k, default in ENTRY_DEFAULTS.items():
                if k not in e:
                    e[k] = default.copy() if isinstance(default, list) else default
    return raw


def _ensure_source(catalog: dict, source_file: str, citation_key: str) -> dict:
    if source_file not in catalog["sources"]:
        catalog["sources"][source_file] = {"citation_key": str(citation_key), "figures": []}
    return catalog["sources"][source_file]


def all_figures(catalog: dict):
    for src_data in catalog["sources"].values():
        yield from src_data["figures"]


def next_fig_id(catalog: dict) -> str:
    existing = [int(e["id"].split("_")[1])
                for e in all_figures(catalog)
                if e.get("id", "").startswith("fig_") and "_" in e["id"]]
    return f"fig_{max(existing, default=0) + 1:03d}"


def already_in_catalog(catalog: dict, source_file: str, path_str: str) -> bool:
    src_data = catalog["sources"].get(source_file)
    if not src_data:
        return False
    return any(e.get("path") == path_str for e in src_data["figures"])


def make_entry(fig_id: str, page: int, path_str: str, needs_crop: bool = False) -> dict:
    entry: dict = {}
    for k, default in ENTRY_DEFAULTS.items():
        entry[k] = default.copy() if isinstance(default, (list, dict)) else default
    entry["id"] = fig_id
    entry["page"] = page
    entry["path"] = path_str
    entry["needs_crop"] = needs_crop
    return entry


def _compute_status(entry: dict) -> str:
    caption_ok = bool(entry.get("caption_verified"))
    has_meta   = bool(entry.get("visual_content"))
    if caption_ok and has_meta:   return "complete"
    if caption_ok:                return "caption-ok"
    if has_meta:                  return "draft"
    return "un-processed"


def _refresh_statuses(catalog: dict) -> None:
    for e in all_figures(catalog):
        e["_status"] = _compute_status(e)


def save_catalog(catalog: dict, path: Path, dry_run: bool = False) -> None:
    if dry_run:
        return
    _refresh_statuses(catalog)
    catalog["_meta"]["last_updated"] = date.today().isoformat()
    path.write_text(json.dumps(catalog, ensure_ascii=False, indent=2), encoding="utf-8")
    guide = path.parent / "CATALOG_GUIDE.md"
    if not guide.exists():
        guide.write_text(CATALOG_GUIDE_TEMPLATE, encoding="utf-8")


def _img_name(page_num: int, fig_idx: int) -> str:
    return f"p{page_num:03d}_fig{fig_idx:03d}.png"


def _rel_path(src: Path, img_name: str) -> str:
    return f"2_clean_inputs/{src.stem}/images/{img_name}"


# ── Language detection ─────────────────────────────────────────────────────────

def detect_lang(src_stem: str, override: dict[str, str]) -> str:
    if src_stem in override:
        return override[src_stem]
    name = src_stem.lower()
    for hint, lang in LANG_HINT.items():
        if hint in name:
            return lang
    return "en"


# ── Keywords auto-derive ──────────────────────────────────────────────────────

def _tokenize(text: str) -> list[str]:
    if not text:
        return []
    words = []
    for w in text.replace(",", " ").replace(".", " ").replace(";", " ").split():
        w = w.strip("()[]{}\"':/\\-_").lower()
        if len(w) >= 3 and w not in STOPWORDS and w.isalpha():
            words.append(w)
    return words


def _auto_keywords(caption: str | None, section_heading: str | None,
                   text_context: str | None) -> list[str] | None:
    tokens: list[str] = []
    if caption:
        cap_tok = _tokenize(caption)
        tokens.extend(cap_tok * 2)   # dupla súly
    if section_heading:
        tokens.extend(_tokenize(section_heading))
    if text_context:
        tokens.extend(_tokenize(text_context[:300]))
    if not tokens:
        return None
    freq = Counter(tokens)
    keywords = [w for w, _ in freq.most_common(12) if len(w) >= 3][:8]
    return keywords if len(keywords) >= 2 else None


# ── MinerU runner ──────────────────────────────────────────────────────────────

# MinerU 2.7.6 backend-ek (03_run_mineru_pipeline.py mintájára)
MINERU_BACKENDS = ["pipeline", "vlm-auto-engine", "vlm-http-client",
                   "hybrid-auto-engine", "hybrid-http-client"]
MINERU_DEFAULT_BACKEND = "pipeline"   # CPU-kompatibilis, GPU nélkül is fut

# ── Worker / VRAM auto-detect ─────────────────────────────────────────────────

def _auto_workers() -> int:
    """Elérhető CPU magok felének fele, minimum 1."""
    return max(1, (os.cpu_count() or 2) // 2)


def _detect_vram_gb(device: str | None) -> int | None:
    """nvidia-smi lekérdezés: adott GPU teljes VRAM GB-ban, vagy None ha nem elérhető."""
    if not device or "cuda" not in device:
        return None
    gpu_idx = 0
    if ":" in device:
        try:
            gpu_idx = int(device.split(":")[1])
        except ValueError:
            pass
    try:
        result = subprocess.run(
            ["nvidia-smi", "--query-gpu=memory.total", "--format=csv,noheader,nounits"],
            capture_output=True, text=True, timeout=5,
        )
        if result.returncode == 0:
            lines = result.stdout.strip().splitlines()
            total_mb = int(lines[min(gpu_idx, len(lines) - 1)])
            return max(1, total_mb // 1024)
    except Exception:
        pass
    return None


def _run_mineru(src: Path, out_parent: Path, lang: str, dry_run: bool,
                backend: str = MINERU_DEFAULT_BACKEND,
                device: str | None = None,
                vram: int | None = None) -> bool:
    """MinerU futtatása egy forrásra (thread-safe: minden hívás külön subprocess).

    backend: MinerU -b flag. 'pipeline' = CPU (default), 'vlm-auto-engine' = GPU.
    device:  MinerU -d flag (csak pipeline). pl. 'cpu', 'cuda', 'cuda:0'.
    vram:    MinerU --vram flag GB-ban (csak pipeline, párhuzamos futásnál GPU split).
    """
    if dry_run:
        extras = ""
        if device and backend == "pipeline":
            extras += f" -d {device}"
        if vram and backend == "pipeline":
            extras += f" --vram {vram}"
        print(f"  [DRY] mineru -p {src.name} -b {backend} -l {lang}{extras}")
        return True
    out_parent.mkdir(parents=True, exist_ok=True)
    cmd = [
        "conda", "run", "-n", MINERU_ENV, "--no-capture-output",
        "mineru", "-p", str(src), "-o", str(out_parent),
        "-m", "auto", "-b", backend, "-l", lang,
    ]
    if device and backend == "pipeline":
        cmd += ["-d", device]
    if vram and backend == "pipeline":
        cmd += ["--vram", str(vram)]
    proc = subprocess.run(cmd, timeout=3600)
    return proc.returncode == 0


def _flatten_mineru(out_parent: Path, src_stem: str) -> Path | None:
    """<out_parent>/<stem>/auto/* → <out_parent>/mineru/"""
    nested = out_parent / src_stem / "auto"
    if not nested.exists():
        nested = out_parent / "auto"
        if not nested.exists():
            return None
    mineru_dir = out_parent / "mineru"
    if mineru_dir.exists():
        shutil.rmtree(mineru_dir)
    mineru_dir.mkdir(parents=True, exist_ok=True)
    for item in nested.iterdir():
        dst = mineru_dir / item.name
        if dst.exists():
            shutil.rmtree(dst) if dst.is_dir() else dst.unlink()
        shutil.move(str(item), str(dst))
    nested.rmdir()
    try:
        (out_parent / src_stem).rmdir()
    except OSError:
        pass
    return mineru_dir


# ── content_list.json parser ──────────────────────────────────────────────────

def _load_content_list(mineru_dir: Path, src_stem: str) -> list[dict] | None:
    cl_path = mineru_dir / f"{src_stem}_content_list.json"
    if not cl_path.exists():
        return None
    try:
        return json.loads(cl_path.read_text(encoding="utf-8"))
    except Exception as e:
        print(f"  ⚠️  content_list.json hiba: {e}", file=sys.stderr)
        return None


def _get_context_around(entries: list[dict], idx: int, window: int) -> str:
    parts = []
    for i in range(max(0, idx - window), min(len(entries), idx + window + 1)):
        if i == idx:
            continue
        e = entries[i]
        if e.get("type") in ("text", "title") and e.get("text"):
            parts.append(e["text"].strip())
    return " ".join(parts)


def _get_section_heading(entries: list[dict], idx: int) -> str | None:
    for i in range(idx - 1, max(0, idx - 30) - 1, -1):
        e = entries[i]
        if e.get("type") in ("text", "title") and e.get("text_level") == 1:
            return e["text"].strip()
    return None


def _process_content_list(
    content_list: list[dict], src: Path, out_dir: Path,
    mineru_images_dir: Path, catalog: dict,
    citation_key: str, dry_run: bool,
) -> tuple[int, int]:
    saved = skipped = 0
    page_counters: dict[int, int] = {}

    for idx, entry in enumerate(content_list):
        if entry.get("type") != "image":
            continue
        img_path_rel = entry.get("img_path", "")
        if not img_path_rel:
            skipped += 1
            continue

        src_img = mineru_images_dir / Path(img_path_rel).name
        if not src_img.exists():
            src_img2 = mineru_images_dir.parent / img_path_rel
            if src_img2.exists():
                src_img = src_img2
            else:
                print(f"  ⚠️  Kép nem található: {img_path_rel}")
                skipped += 1
                continue

        page_num = int(entry.get("page_idx", 0)) + 1
        page_counters[page_num] = page_counters.get(page_num, 0) + 1
        fig_idx  = page_counters[page_num]
        img_name = _img_name(page_num, fig_idx)
        rel_path = _rel_path(src, img_name)

        if already_in_catalog(catalog, src.name, rel_path):
            skipped += 1
            continue

        if not dry_run:
            dst = out_dir / "images" / img_name
            dst.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(str(src_img), str(dst))

        # Szemantikus mezők
        captions_raw = entry.get("image_caption", [])
        caption      = captions_raw[0].strip() if captions_raw else None
        text_context = _get_context_around(content_list, idx, CONTEXT_WINDOW) or None
        heading      = _get_section_heading(content_list, idx)
        # Fallback caption: ha MinerU nem talált image_caption-t (pl. prezentáció-PDF
        # vagy diasor-kép), a legközelebbi section heading átveszi a caption szerepét
        if not caption and heading:
            caption = heading
        keywords     = _auto_keywords(caption, heading, text_context)

        e = make_entry(next_fig_id(catalog), page_num, rel_path, needs_crop=False)
        if caption:      e["caption"]      = caption
        if text_context: e["text_context"] = text_context
        if keywords:     e["keywords"]     = keywords

        _ensure_source(catalog, src.name, citation_key)["figures"].append(e)
        saved += 1

    return saved, skipped


# ── PPTX extractor ─────────────────────────────────────────────────────────────

def _collect_pptx_images(shapes) -> list:
    result = []
    for shape in shapes:
        if shape.shape_type == 6 and hasattr(shape, "shapes"):
            result.extend(_collect_pptx_images(shape.shapes))
            continue
        try:
            from lxml import etree
            xml = etree.tostring(shape.element).decode("utf-8")
        except Exception:
            xml = getattr(shape.element, "xml", "")
        if "blipFill" in xml or "a:blip" in xml:
            result.append(shape)
    return result


def _extract_pptx(src: Path, out_dir: Path, citation_key: str,
                  catalog: dict, dry_run: bool) -> tuple[int, int]:
    try:
        from pptx import Presentation
    except ImportError:
        print("  HIBA: python-pptx nincs telepítve.", file=sys.stderr)
        return 0, 0

    prs = Presentation(str(src))
    saved = skipped = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        # Slide title → caption (python-pptx shapes.title)
        slide_title: str | None = None
        try:
            t = slide.shapes.title
            if t and t.has_text_frame:
                slide_title = t.text.strip() or None
        except Exception:
            pass

        # Slide szöveg (title KIVÉTELÉVEL) → text_context
        title_shape_id = slide.shapes.title.shape_id if slide.shapes.title else None
        slide_text_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.shape_id == title_shape_id:
                continue   # title külön kezeljük
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    slide_text_parts.append(t)
        # Ha nincs body-szöveg, fallback: a title adja a text_context-et is
        slide_body = " ".join(slide_text_parts) or None
        slide_text_context = slide_body or slide_title

        img_shapes  = _collect_pptx_images(slide.shapes)
        page_fig_idx = 0
        for shape in img_shapes:
            try:
                image = shape.image
                blob  = image.blob
                w = shape.width.pt if hasattr(shape.width, "pt") else 0
                h = shape.height.pt if hasattr(shape.height, "pt") else 0
                if w * h < MIN_AREA / 10:
                    skipped += 1
                    continue
                page_fig_idx += 1
                img_name = _img_name(slide_idx, page_fig_idx)
                rel      = _rel_path(src, img_name)

                if already_in_catalog(catalog, src.name, rel):
                    skipped += 1
                    continue

                try:
                    import fitz
                    pix = fitz.Pixmap(blob)
                    if pix.n > 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    img_bytes = pix.tobytes("png")
                except Exception:
                    img_bytes = blob

                if not dry_run:
                    img_path = out_dir / "images" / img_name
                    img_path.parent.mkdir(parents=True, exist_ok=True)
                    img_path.write_bytes(img_bytes)

                e = make_entry(next_fig_id(catalog), slide_idx, rel, needs_crop=False)
                if slide_title:
                    e["caption"] = slide_title
                if slide_text_context:
                    e["text_context"] = slide_text_context[:500]
                    kw = _auto_keywords(slide_title, None, slide_text_context)
                    if kw:
                        e["keywords"] = kw

                _ensure_source(catalog, src.name, citation_key)["figures"].append(e)
                saved += 1

            except Exception as ex:
                print(f"  SKIP slide {slide_idx}: {ex}", file=sys.stderr)

    return saved, skipped


# ── Per-source worker (párhuzamos futáshoz) ────────────────────────────────────

def _merge_tmp_catalog(tmp_catalog: dict, catalog: dict,
                       catalog_lock: threading.Lock) -> None:
    """tmp_catalog["sources"] mergeLése a shared catalog-ba (lock alatt)."""
    with catalog_lock:
        for src_name, src_data in tmp_catalog["sources"].items():
            if src_name not in catalog["sources"]:
                catalog["sources"][src_name] = src_data
            else:
                existing_paths = {e["path"]
                                  for e in catalog["sources"][src_name]["figures"]}
                for e in src_data["figures"]:
                    if e["path"] not in existing_paths:
                        catalog["sources"][src_name]["figures"].append(e)


def _process_source(
    src: Path,
    clean_in: Path,
    citations: dict,
    catalog: dict,
    catalog_lock: threading.Lock,
    override: dict[str, str],
    backend: str,
    device: str | None,
    vram: int | None,
    dry_run: bool,
    resume: bool = True,
) -> tuple[int, int, str]:
    """Egy forrás teljes feldolgozása (MinerU vagy PPTX). Thread-safe.

    resume=True (default): ha <stem>/mineru/<stem>_content_list.json már létezik,
    átugorja a MinerU újrafuttatását — folytatja az előző megszakított futást.
    resume=False: mindig újra futtatja MinerU-t (clean run).

    Visszatér: (saved, skipped, label)
    """
    cit     = citations.get(src.name, "?")
    out_dir = clean_in / src.stem

    if src.suffix.lower() == ".pptx":
        tmp_catalog: dict = {"_meta": {}, "sources": {}}
        s, k = _extract_pptx(src, out_dir, cit, tmp_catalog, dry_run)
        _merge_tmp_catalog(tmp_catalog, catalog, catalog_lock)
        return s, k, "PPTX (python-pptx)"

    # PDF: MinerU
    lang = detect_lang(src.stem, override)

    # Resume: ha content_list.json már létezik, átugorjuk MinerU újrafuttatását
    existing_cl = out_dir / "mineru" / f"{src.stem}_content_list.json"
    if resume and existing_cl.exists():
        mineru_dir   = out_dir / "mineru"
        content_list = _load_content_list(mineru_dir, src.stem)
        if content_list is not None:
            tmp_catalog = {"_meta": {}, "sources": {}}
            s, k = _process_content_list(
                content_list, src, out_dir,
                mineru_dir / "images", tmp_catalog, cit, dry_run)
            _merge_tmp_catalog(tmp_catalog, catalog, catalog_lock)
            return s, k, f"MinerU ({lang}, RESUMED — content_list kész volt)"

    # Normál MinerU futás
    ok = _run_mineru(src, out_dir, lang, dry_run,
                     backend=backend, device=device, vram=vram)
    if not ok:
        return 0, 0, "✗ MinerU hiba"

    if dry_run:
        return 0, 0, f"MinerU ({lang}, {backend})"

    mineru_dir = _flatten_mineru(out_dir, src.stem)
    if not mineru_dir:
        return 0, 0, "⚠️  MinerU output hiányzik"

    content_list = _load_content_list(mineru_dir, src.stem)
    if content_list is None:
        return 0, 0, "⚠️  content_list.json hiányzik"

    tmp_catalog = {"_meta": {}, "sources": {}}
    s, k = _process_content_list(
        content_list, src, out_dir,
        mineru_dir / "images", tmp_catalog, cit, dry_run)
    _merge_tmp_catalog(tmp_catalog, catalog, catalog_lock)
    return s, k, f"MinerU ({lang}, {backend})"


# ── Főfüggvény ─────────────────────────────────────────────────────────────────

def main() -> int:
    parser = argparse.ArgumentParser(
        description="MinerU-first ábra-kinyerő + katalógus-builder (v4 séma)")
    parser.add_argument("--week-dir", required=True, type=Path)
    parser.add_argument("--source",   type=str, default=None)
    parser.add_argument("--lang-map", type=str, default="",
                        help="stem=lang,... pl. 'nagyi2013_eloadas=latin'")
    parser.add_argument("--backend", default=MINERU_DEFAULT_BACKEND,
                        choices=MINERU_BACKENDS,
                        help=f"MinerU backend (default: {MINERU_DEFAULT_BACKEND}). "
                             "'pipeline'=CPU; 'vlm-auto-engine'=GPU lokálisan.")
    parser.add_argument("--device", type=str, default=None,
                        help="MinerU -d flag, csak pipeline backendhez. "
                             "pl. 'cpu', 'cuda', 'cuda:0'. Default: auto-detect.")
    parser.add_argument("--workers", type=int, default=None,
                        help="Párhuzamos worker-ek száma (default: cpu_count // 2). "
                             "1 = soros feldolgozás.")
    parser.add_argument("--vram-per-worker", type=int, default=None, dest="vram_per_worker",
                        help="GPU VRAM GB/worker (csak pipeline backend). "
                             "Default: auto (nvidia-smi / workers).")
    parser.add_argument("--no-resume", action="store_true",
                        help="Mindig újra futtatja MinerU-t, még ha content_list.json "
                             "már létezik is (clean run). Default: resume=True.")
    parser.add_argument("--dry-run",  action="store_true")
    args = parser.parse_args()

    week_dir = args.week_dir.resolve()
    raw_in   = week_dir / "1_raw_inputs"
    clean_in = week_dir / "2_clean_inputs"
    cat_path = clean_in / "figure_catalog.json"

    if not raw_in.is_dir():
        sys.exit(f"HIBA: nem találhato {raw_in}")
    clean_in.mkdir(parents=True, exist_ok=True)

    override: dict[str, str] = {}
    for chunk in (args.lang_map or "").split(","):
        if "=" in chunk:
            s, l = chunk.strip().split("=", 1)
            override[s.strip()] = l.strip()

    catalog       = load_catalog(cat_path)
    catalog_lock  = threading.Lock()
    citations     = load_citations(week_dir)
    prefix        = "[DRY] " if args.dry_run else ""

    if args.source:
        sources = [raw_in / args.source]
        if not sources[0].exists():
            sys.exit(f"HIBA: nem található {sources[0]}")
    else:
        sources = sorted(f for f in raw_in.iterdir()
                         if f.is_file()
                         and f.suffix.lower() in (".pdf", ".pptx")
                         and not f.name.startswith("_")
                         and f.name != "citations.json")

    # Worker szám + VRAM meghatározása
    workers = args.workers if args.workers else _auto_workers()
    workers = min(workers, len(sources))  # felesleges worker-eket ne indítson

    vram: int | None = args.vram_per_worker
    if vram is None and args.device and "cuda" in args.device:
        total_vram = _detect_vram_gb(args.device)
        if total_vram:
            vram = max(1, total_vram // workers)

    cpu_info = f"cpu_count={os.cpu_count()}, workers={workers}"
    vram_info = f", vram/worker={vram}GB" if vram else ""
    print(f"{prefix}Forrás-feldolgozás: {len(sources)} forrás | {cpu_info}{vram_info}"
          f" | backend={args.backend}")

    total_saved = total_skip = 0

    with concurrent.futures.ThreadPoolExecutor(max_workers=workers) as executor:
        future_to_src = {
            executor.submit(
                _process_source,
                src, clean_in, citations, catalog, catalog_lock,
                override, args.backend, args.device, vram, args.dry_run,
                resume=not args.no_resume,
            ): src
            for src in sources
        }
        for future in concurrent.futures.as_completed(future_to_src):
            src = future_to_src[future]
            try:
                s, k, label = future.result()
            except Exception as exc:
                print(f"  ✗ {src.name}: {exc}", file=sys.stderr)
                s = k = 0
                label = "HIBA"
            print(f"  ✓ {src.name} [{label}] → {s} kép, {k} skip")
            total_saved += s
            total_skip  += k

    save_catalog(catalog, cat_path, args.dry_run)

    total = sum(1 for _ in all_figures(catalog))
    cap_n = sum(1 for e in all_figures(catalog) if e.get("caption"))
    ctx_n = sum(1 for e in all_figures(catalog) if e.get("text_context"))
    kw_n  = sum(1 for e in all_figures(catalog) if e.get("keywords"))
    stat  = Counter(e["_status"] for e in all_figures(catalog))

    print(f"\n{prefix}Kész: {total_saved} kép | {total_skip} skip")
    print(f"Katalógus: {total} entry | {dict(stat)}")
    if total:
        print(f"Kitöltöttség: caption {cap_n}/{total} ({cap_n/total:.0%})"
              f" | text_context {ctx_n}/{total} ({ctx_n/total:.0%})"
              f" | keywords {kw_n}/{total} ({kw_n/total:.0%})")
    return 0


if __name__ == "__main__":
    sys.exit(main())
