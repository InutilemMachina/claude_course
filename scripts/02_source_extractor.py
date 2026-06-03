"""
02_source_extractor.py -- Ábra-kinyerő a 1_raw_inputs/ forrásokból.

Minden forrástípusból PNG képeket ment 2_clean_inputs/<stem>/images/-ba,
és felépíti/frissíti a 2_clean_inputs/figure_catalog.json-t.

    PDF (born-digital): beágyazott képek kinyerése; kicsi dekorációk kihagyva.
    PDF (szkennelt):    teljes oldal renderelése PNG-ként, needs_crop: true + figyelmeztetés.
    PPTX:               diák beágyazott képeinek kinyerése.

A szöveg-szintézist Claude végzi közvetlenül a forrásból — itt CSAK ábrák kellenek.

Usage:
    python scripts/02_source_extractor.py --week-dir test_outputs/atg/1_het
    python scripts/02_source_extractor.py --week-dir test_outputs/atg/1_het --dry-run
"""

import argparse
import importlib.util
import json
import sys
from pathlib import Path

try:
    from _encoding_fix import fix_stdout as _fix_stdout
    _fix_stdout()
except ImportError:
    pass

# ── Auto-crop betöltése (scripts/_auto_crop.py) ────────────────────────────────
try:
    _ac_spec = importlib.util.spec_from_file_location(
        "_auto_crop", Path(__file__).parent / "_auto_crop.py")
    _ac_mod = importlib.util.module_from_spec(_ac_spec)
    _ac_spec.loader.exec_module(_ac_mod)
    _auto_crop_fn = _ac_mod.auto_crop
except Exception:
    _auto_crop_fn = None  # Pillow hiánya vagy hiba: auto-crop kikapcsolva

# ── Küszöbértékek ──────────────────────────────────────────────────────────────
MIN_AREA   = 10_000   # px² alatt: dekoráció/logó → kihagyva (pl. 100×100)
PAGE_FILL  = 0.85     # oldal-terület %-a felett: szkennelt oldal → render + warn
RENDER_DPI = 150      # szkennelt oldal renderelési felbontása


# ── Segédfüggvények ────────────────────────────────────────────────────────────

def load_citations(week_dir: Path) -> dict:
    """citations.json betöltése fájlnév→citáció-kulcs mappinghoz."""
    cit_path = week_dir / "1_raw_inputs" / "citations.json"
    if not cit_path.exists():
        return {}
    data = json.loads(cit_path.read_text(encoding="utf-8"))
    return {v["filename"]: k for k, v in data.items()
            if k != "_meta" and v.get("filename")}


def next_fig_id(catalog: list) -> str:
    """Következő fig_NNN azonosító a katalógus alapján."""
    existing = [int(e["id"].split("_")[1]) for e in catalog if "_" in e.get("id","")]
    n = max(existing, default=0) + 1
    return f"fig_{n:03d}"


def already_in_catalog(catalog: list, source_file: str, filename: str) -> bool:
    """Ellenőrzi, hogy a (source_file, filename) pár már szerepel-e."""
    return any(e.get("source_file") == source_file and e.get("filename") == filename
               for e in catalog)


def save_catalog(catalog: list, path: Path, dry_run: bool):
    if not dry_run:
        path.write_text(json.dumps(catalog, ensure_ascii=False, indent=2), encoding="utf-8")


# ── PDF extractor ──────────────────────────────────────────────────────────────

def extract_pdf(src: Path, out_dir: Path, citation_key: str,
                catalog: list, dry_run: bool,
                specific_pages: list[int] | None = None) -> tuple[int, int, int]:
    """
    Born-digital: beágyazott képek kinyerése.
    Vegyes PDF (néhány szkennelt oldal, <SCANNED_THRESHOLD): crop-ra vár + figyelmeztetés.
    Teljes szkennelt PDF (>=SCANNED_THRESHOLD): kihagyva, 1 összesítő figyelmeztetés.
    specific_pages: ha megadott (Claude azonosította), csak ezeket az oldalakat dolgozza fel —
                    oldalanként annyi PNG-t ment, ahány kép van rajta (szkenneltnél 1 oldal-render).
    Visszatérés: (mentett, kihagyott_deko, crop_figyelmeztetések)
    """
    SCANNED_THRESHOLD = 0.50   # ha >50% oldal szkennelt → egész forrás kihagyva
    try:
        import fitz
    except ImportError:
        print(f"  HIBA: PyMuPDF nincs telepítve. pip install pymupdf", file=sys.stderr)
        return 0, 0, 0

    doc = fitz.open(str(src))
    n_pages = len(doc)

    # specific_pages mód: Claude azonosított oldalak feldolgozása.
    # Az ismétlések adják meg a képszámot: "5,12,12,12" → p5: 1 kép, p12: 3 kép.
    # Szkennelt oldalnál: N kép → N külön fájl (p012_fig001.png, p012_fig002.png, …)
    # Born-digital oldalnál: a beágyazott képek alapján mentünk, ismétlések figyelmen kívül.
    if specific_pages is not None:
        from collections import Counter
        page_counts = Counter(specific_pages)  # {oldal: hány képet kér}
        saved = skipped = crop_warn = 0

        for page_num, fig_count in sorted(page_counts.items()):
            if page_num < 1 or page_num > n_pages:
                print(f"  SKIP  oldal {page_num}: kívül esik ({n_pages} old.)")
                continue
            page = doc[page_num - 1]
            page_area = page.rect.width * page.rect.height
            page_imgs = page.get_images(full=True)

            # Szkennelt oldal detektálás (első kép > PAGE_FILL?)
            is_scanned = any(
                (doc.extract_image(img[0])["width"] * doc.extract_image(img[0])["height"])
                / page_area >= PAGE_FILL
                for img in page_imgs[:1]  # csak az első képet nézzük (gyors)
            ) if page_imgs else False

            if is_scanned:
                # Szkennelt oldal: N = fig_count másolatot ment, külön fájlnévvel
                mat = fitz.Matrix(RENDER_DPI / 72, RENDER_DPI / 72)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                for fig_idx in range(1, fig_count + 1):
                    img_name = f"p{page_num:03d}_fig{fig_idx:03d}.png"
                    rel_path = f"2_clean_inputs/{src.stem}/images/{img_name}"
                    print(f"  ⚠️  CROP SZÜKSÉGES: oldal {page_num} kép {fig_idx} → {img_name}")
                    if not dry_run:
                        img_path = out_dir / "images" / img_name
                        img_path.parent.mkdir(parents=True, exist_ok=True)
                        img_path.write_bytes(img_bytes)
                    if not already_in_catalog(catalog, src.name, rel_path):
                        catalog.append({
                            "id": next_fig_id(catalog),
                            "source_file": src.name,
                            "citation_key": citation_key,
                            "page": page_num,
                            "filename": rel_path,
                            "needs_crop": True,
                            "caption": None,
                            "suggested_section": None,
                        })
                        saved += 1
                        crop_warn += 1
                    else:
                        skipped += 1
            else:
                # Born-digital oldal: beágyazott képek alapján mentünk.
                # Ha nincs raszterkép (csak vektoros ábrák), oldalrenderelés → needs_crop.
                raster_saved = 0
                for img_idx, img_info in enumerate(page_imgs, 1):
                    xref = img_info[0]
                    base = doc.extract_image(xref)
                    w, h = base["width"], base["height"]
                    if w * h < MIN_AREA:
                        skipped += 1
                        continue
                    img_name = f"p{page_num:03d}_img{img_idx:03d}.png"
                    img_bytes = base["image"]
                    if base["ext"] != "png":
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n > 4:
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                        img_bytes = pix.tobytes("png")
                    rel_path = f"2_clean_inputs/{src.stem}/images/{img_name}"
                    if not dry_run:
                        img_path = out_dir / "images" / img_name
                        img_path.parent.mkdir(parents=True, exist_ok=True)
                        img_path.write_bytes(img_bytes)
                    if not already_in_catalog(catalog, src.name, rel_path):
                        catalog.append({
                            "id": next_fig_id(catalog),
                            "source_file": src.name,
                            "citation_key": citation_key,
                            "page": page_num,
                            "filename": rel_path,
                            "needs_crop": False,
                            "caption": None,
                            "suggested_section": None,
                        })
                        saved += 1
                        raster_saved += 1
                    else:
                        skipped += 1

                has_raster = any(
                    (doc.extract_image(img[0])["width"] * doc.extract_image(img[0])["height"]) >= MIN_AREA
                    for img in page_imgs
                )
                if not has_raster:
                    # Nincs raszterkép → vektoros ábra: oldalrenderelés + auto-crop kísérlet
                    mat = fitz.Matrix(RENDER_DPI / 72, RENDER_DPI / 72)
                    pix = page.get_pixmap(matrix=mat)
                    img_bytes = pix.tobytes("png")
                    for fig_idx in range(1, fig_count + 1):
                        img_name = f"p{page_num:03d}_fig{fig_idx:03d}.png"
                        rel_path = f"2_clean_inputs/{src.stem}/images/{img_name}"
                        needs_crop = True
                        if not dry_run:
                            img_path = out_dir / "images" / img_name
                            img_path.parent.mkdir(parents=True, exist_ok=True)
                            img_path.write_bytes(img_bytes)
                            if _auto_crop_fn is not None:
                                cropped, ratio = _auto_crop_fn(img_path)
                                if cropped:
                                    needs_crop = False
                                    print(f"  ✂️  AUTO-CROP: oldal {page_num} kép {fig_idx} → {img_name} ({ratio:.0%} levágva)")
                                else:
                                    print(f"  ⚠️  CROP SZÜKSÉGES (vektor): oldal {page_num} kép {fig_idx} → {img_name}")
                            else:
                                print(f"  ⚠️  CROP SZÜKSÉGES (vektor): oldal {page_num} kép {fig_idx} → {img_name}")
                        else:
                            print(f"  ⚠️  CROP SZÜKSÉGES (vektor): oldal {page_num} kép {fig_idx} → {img_name}")
                        if not already_in_catalog(catalog, src.name, rel_path):
                            catalog.append({
                                "id": next_fig_id(catalog),
                                "source_file": src.name,
                                "citation_key": citation_key,
                                "page": page_num,
                                "filename": rel_path,
                                "needs_crop": needs_crop,
                                "caption": None,
                                "suggested_section": None,
                            })
                            saved += 1
                            if needs_crop:
                                crop_warn += 1
                        else:
                            # Meglévő bejegyzés: ha auto-crop lefutott → frissítjük needs_crop-ot
                            if not needs_crop:
                                for e in catalog:
                                    if e.get("source_file") == src.name and e.get("filename") == rel_path:
                                        e["needs_crop"] = False
                                        break
                            skipped += 1

        doc.close()
        return saved, skipped, crop_warn

    # 1. átmenet: szkennelt oldalak arányának meghatározása
    scanned_pages = []
    for page_num, page in enumerate(doc, 1):
        page_area = page.rect.width * page.rect.height
        for img_info in page.get_images(full=True):
            base = doc.extract_image(img_info[0])
            if (base["width"] * base["height"]) / page_area >= PAGE_FILL:
                scanned_pages.append(page_num)
                break  # oldalanként elég egy találat

    scanned_ratio = len(scanned_pages) / n_pages if n_pages else 0

    # Teljes szkennelt dokumentum → kihagyás
    if scanned_ratio >= SCANNED_THRESHOLD:
        print(f"  ⚠️  SZKENNELT FORRÁS ({len(scanned_pages)}/{n_pages} oldal, "
              f"{scanned_ratio:.0%}) — ábra-kinyerés kihagyva. "
              f"Használd: --source {src.name} --pages <oldalszámok> (Claude azonosítja)")
        doc.close()
        return 0, 0, 0

    # 2. átmenet: tényleges feldolgozás (vegyes vagy born-digital)
    saved = skipped = crop_warn = 0

    for page_num, page in enumerate(doc, 1):
        page_area = page.rect.width * page.rect.height
        is_scanned_page = page_num in scanned_pages

        for img_idx, img_info in enumerate(page.get_images(full=True), 1):
            xref = img_info[0]
            base = doc.extract_image(xref)
            w, h = base["width"], base["height"]
            area = w * h

            # Dekoráció/logó szűrő
            if area < MIN_AREA:
                skipped += 1
                continue

            needs_crop = False
            if is_scanned_page:
                # Vegyes PDF: szkennelt oldal renderelése + crop figyelmeztetés
                mat = fitz.Matrix(RENDER_DPI / 72, RENDER_DPI / 72)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                img_name = f"p{page_num:03d}_page.png"
                needs_crop = True
                crop_warn += 1
                print(f"  ⚠️  CROP SZÜKSÉGES: oldal {page_num} → {img_name}")
            else:
                img_name = f"p{page_num:03d}_img{img_idx:03d}.png"
                img_bytes = base["image"]
                if base["ext"] != "png":
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n > 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    img_bytes = pix.tobytes("png")

            rel_path = f"2_clean_inputs/{src.stem}/images/{img_name}"

            if not dry_run:
                img_path = out_dir / "images" / img_name
                img_path.parent.mkdir(parents=True, exist_ok=True)
                img_path.write_bytes(img_bytes)

            if not already_in_catalog(catalog, src.name, rel_path):
                catalog.append({
                    "id": next_fig_id(catalog),
                    "source_file": src.name,
                    "citation_key": citation_key,
                    "page": page_num,
                    "filename": rel_path,
                    "needs_crop": needs_crop,
                    "caption": None,
                    "suggested_section": None,
                })
            saved += 1

    doc.close()
    return saved, skipped, crop_warn


# ── PPTX extractor ─────────────────────────────────────────────────────────────

def _collect_pptx_images(shapes) -> list:
    """Minden blipFill-tartalmú shape képét gyűjti rekurzívan (group-on belül is)."""
    result = []
    for shape in shapes:
        # Rekurzív: group shape
        if shape.shape_type == 6 and hasattr(shape, "shapes"):
            result.extend(_collect_pptx_images(shape.shapes))
            continue
        # XML-alapú detektálás — shape típustól független
        try:
            from lxml import etree
            xml = etree.tostring(shape.element).decode("utf-8")
        except Exception:
            xml = getattr(shape.element, "xml", "")
        if "blipFill" not in xml and "a:blip" not in xml:
            continue
        try:
            result.append(shape)
        except Exception:
            pass
    return result


def extract_pptx(src: Path, out_dir: Path, citation_key: str,
                 catalog: list, dry_run: bool) -> tuple[int, int]:
    """Diák beágyazott képeinek kinyerése PNG-ként (minden shape-típusból)."""
    try:
        from pptx import Presentation
    except ImportError:
        print(f"  HIBA: python-pptx nincs telepítve. pip install python-pptx", file=sys.stderr)
        return 0, 0

    prs = Presentation(str(src))
    saved = skipped = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        img_shapes = _collect_pptx_images(slide.shapes)
        for img_idx, shape in enumerate(img_shapes, 1):
            try:
                image = shape.image
                blob = image.blob
                w = shape.width.pt if hasattr(shape.width, "pt") else 0
                h = shape.height.pt if hasattr(shape.height, "pt") else 0
                area = w * h

                if area < MIN_AREA / 10:  # PPTX pt vs px – lazább küszöb
                    skipped += 1
                    continue

                img_name = f"slide{slide_idx:03d}_img{img_idx:03d}.png"
                rel_path = f"2_clean_inputs/{src.stem}/images/{img_name}"

                # PNG-vé konvertálás (fitz-cel ha elérhető, különben nyers)
                try:
                    import fitz
                    pix = fitz.Pixmap(blob)
                    if pix.n > 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    img_bytes = pix.tobytes("png")
                except Exception:
                    img_bytes = blob  # eredeti formátum fallback

                if not dry_run:
                    img_path = out_dir / "images" / img_name
                    img_path.parent.mkdir(parents=True, exist_ok=True)
                    img_path.write_bytes(img_bytes)

                if not already_in_catalog(catalog, src.name, rel_path):
                    catalog.append({
                        "id": next_fig_id(catalog),
                        "source_file": src.name,
                        "citation_key": citation_key,
                        "page": slide_idx,
                        "filename": rel_path,
                        "needs_crop": False,
                        "caption": None,
                        "suggested_section": None,
                    })
                    saved += 1
            except Exception as e:
                print(f"  SKIP  slide {slide_idx} kep: {e}", file=sys.stderr)

    return saved, skipped


# ── Főfüggvény ─────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Ábra-kinyerő: PDF/PPTX → PNG + figure_catalog.json"
    )
    parser.add_argument("--week-dir", required=True, type=Path,
                        help="Heti mappa (pl. test_outputs/atg/1_het)")
    parser.add_argument("--source", type=str, default=None,
                        help="Forrás fájlneve (csak ehhez futtat, --pages-szel)")
    parser.add_argument("--pages", type=str, default=None,
                        help="Claude-azonosított oldalszámok vesszővel (pl. '5,12,23')")
    parser.add_argument("--dry-run", action="store_true",
                        help="Listázás mentés nélkül")
    parser.add_argument("--sync-crop-tasks", action="store_true",
                        help="[x] jelölések szinkronizálása a catalog-ba, majd kilép")
    args = parser.parse_args()

    week_dir  = args.week_dir.resolve()
    raw_in    = week_dir / "1_raw_inputs"
    clean_in  = week_dir / "2_clean_inputs"
    cat_path  = clean_in / "figure_catalog.json"

    # --sync-crop-tasks: [x] bejegyzések catalog-ba írása, majd kilép
    if args.sync_crop_tasks:
        try:
            _ct_spec = importlib.util.spec_from_file_location(
                "_crop_tasks", Path(__file__).parent / "_crop_tasks.py")
            _ct_mod = importlib.util.module_from_spec(_ct_spec)
            _ct_spec.loader.exec_module(_ct_mod)
            n = _ct_mod.sync_crop_tasks(week_dir)
            print(f"Sync kész: {n} bejegyzés frissítve.")
        except Exception as e:
            print(f"HIBA: sync_crop_tasks sikertelen: {e}", file=sys.stderr)
            sys.exit(1)
        sys.exit(0)

    if not raw_in.is_dir():
        sys.exit(f"HIBA: nem található {raw_in}")

    # Meglévő katalógus betöltése (idempotens futtatáshoz)
    catalog = json.loads(cat_path.read_text(encoding="utf-8")) if cat_path.exists() else []
    already = {e["filename"] for e in catalog}

    citations = load_citations(week_dir)
    prefix    = "[DRY] " if args.dry_run else ""

    # --source + --pages: Claude-azonosított szkennelt oldalak feldolgozása
    specific_pages = None
    if args.source and args.pages:
        try:
            specific_pages = [int(p.strip()) for p in args.pages.split(",") if p.strip()]
        except ValueError:
            sys.exit("HIBA: --pages csak egész számokat fogad el, vesszővel (pl. '5,12,23')")

    total_saved = total_skip = total_crop = 0
    processed = []

    # Ha --source megadott: csak azt a fájlt dolgozzuk fel
    sources = []
    if args.source:
        src = raw_in / args.source
        if not src.exists():
            sys.exit(f"HIBA: nem található {src}")
        sources = [src]
    else:
        sources = sorted(f for f in raw_in.iterdir()
                         if f.is_file()
                         and f.suffix.lower() in (".pdf", ".pptx")
                         and not f.name.startswith("_")
                         and f.name != "citations.json")

    for src in sources:
        ext = src.suffix.lower()
        citation_key = citations.get(src.name, "?")
        out_dir = clean_in / src.stem

        label = f"oldal {specific_pages}" if specific_pages else ""
        print(f"{prefix}{src.name} [cit:{citation_key}]{' ' + label if label else ''}")

        if ext == ".pdf":
            saved, skip, crop = extract_pdf(src, out_dir, citation_key,
                                            catalog, args.dry_run, specific_pages)
            total_crop += crop
        elif ext == ".pptx":
            saved, skip = extract_pptx(src, out_dir, citation_key,
                                       catalog, args.dry_run)
            crop = 0
        else:
            continue

        total_saved += saved
        total_skip  += skip
        processed.append(src.name)
        print(f"  → {saved} kép mentve, {skip} dekoráció kihagyva"
              + (f", {crop} oldal crop-ra vár ⚠️" if crop else ""))

    save_catalog(catalog, cat_path, args.dry_run)

    # _crop_tasks.md generálás (csak ha nem dry-run)
    if not args.dry_run:
        try:
            _ct_spec = importlib.util.spec_from_file_location(
                "_crop_tasks", Path(__file__).parent / "_crop_tasks.py")
            _ct_mod = importlib.util.module_from_spec(_ct_spec)
            _ct_spec.loader.exec_module(_ct_mod)
            _ct_mod.generate_crop_tasks(week_dir)
        except Exception as e:
            print(f"  WARN: _crop_tasks.md generálás sikertelen: {e}", file=sys.stderr)

    print(f"\n{prefix}Kész: {len(processed)} forrás | "
          f"{total_saved} kép | {total_skip} deko kihagyva"
          + (f" | ⚠️  {total_crop} oldal CROP SZÜKSÉGES" if total_crop else ""))

    if total_crop:
        print("  Crop-ra váró oldalak a figure_catalog.json-ban: needs_crop: true")


if __name__ == "__main__":
    main()
