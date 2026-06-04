"""
02_image_extraction.py — Ábra-kinyerő a 1_raw_inputs/ forrásokból (v4 séma).

Minden forrásból PNG képeket ment 2_clean_inputs/<stem>/images/-ba egységes
pNNN_figNNN.png néven, és felépíti/frissíti a 2_clean_inputs/figure_catalog.json-t
(v4 séma: _meta + sources csoportosítva).

  PDF (born-digital):       beágyazott + vektoros oldal-render, auto-crop kísérlet
  PDF (szkennelt):          teljes oldal renderelve, needs_crop:true + OCR cache
  PPTX:                     diák beágyazott képei (slide N = page N)

OCR: szkennelt oldalakhoz pytesseract (opcionális) → text/pNNN.txt cache.
A 02b_figure_enricher skill ezt fogyasztja text_context feltöltéshez.

Usage:
    python scripts/02_image_extraction.py --week-dir test_outputs/atg/1_het
    python scripts/02_image_extraction.py --week-dir test_outputs/atg/1_het --source X.pdf --pages "5,12,23"
    python scripts/02_image_extraction.py --week-dir test_outputs/atg/1_het --dry-run
"""

import argparse
import importlib.util
import json
import sys
from collections import Counter
from datetime import date
from pathlib import Path

try:
    from _encoding_fix import fix_stdout as _fix_stdout
    _fix_stdout()
except ImportError:
    pass

# ── Auto-crop betöltése (scripts/_auto_crop.py) ────────────────────────────────
try:
    _ac_spec = importlib.util.spec_from_file_location(
        "_auto_crop", Path(__file__).parent / "_auto_crop.py")
    _ac_mod = importlib.util.module_from_spec(_ac_spec)
    _ac_spec.loader.exec_module(_ac_mod)
    _auto_crop_fn = _ac_mod.auto_crop
except Exception:
    _auto_crop_fn = None  # Pillow hiánya: auto-crop kikapcsolva

# ── pytesseract OCR (opcionális) ───────────────────────────────────────────────
try:
    import pytesseract  # type: ignore
    from PIL import Image  # type: ignore
    _OCR_AVAILABLE = True
except Exception:
    _OCR_AVAILABLE = False

# ── Küszöbértékek ──────────────────────────────────────────────────────────────
MIN_AREA             = 10_000  # px² alatt: dekoráció/logó → kihagyva
PAGE_FILL            = 0.85    # oldal-terület %-a felett: szkennelt oldal
RENDER_DPI           = 150     # szkennelt oldal renderelési felbontása
VECTOR_MIN_DRAW      = 10      # ennyi szignifikáns drawing-elem felett: vektoros ábra
VECTOR_MIN_ELEM_AREA = 200     # pt² — ez alatt: bullet/underline → nem számít
VECTOR_DIVIDER_H     = 5       # pt — alacsonyabb + széles → fejléc-vonal, skip
SCANNED_THRESHOLD    = 0.50    # >50% szkennelt → teljes forrás skip
OCR_LANGS            = "eng+hun"

# ── v4 séma ────────────────────────────────────────────────────────────────────
SCHEMA_VERSION = 4

# Mezők alapértékei az új bejegyzéseknél (LOGIKAI sorrendben — a make_entry
# is ezt a sorrendet írja a dict-be).
ENTRY_DEFAULTS = {
    # 1. Identitás
    "id":               None,    # 🐍 fig_NNN
    "page":             None,    # 🐍 forrás-oldal
    "path":             None,    # 🐍 kép path 2_clean_inputs/.../pNNN_figNNN.png
    # 2. Operatív státusz
    "needs_crop":       False,
    # 3. Ember-olvasható
    "caption":          None,    # 🤖→😎
    "caption_verified": False,   # 😎-only
    # 4. Szemantikus (retrieval — 02b tölti)
    "visual_content":   None,    # 🤖→😎
    "text_context":     None,    # 🤖→😎
    "keywords":         None,    # 🤖→😎  null=un-processed (nem []=no-results)
    # 5. Összegző + user
    "_status":          "un-processed",  # 🐍 derived
    "notes":            [],      # 😎
}

# ── CATALOG_GUIDE.md template (generálódik 2_clean_inputs/ mellé) ─────────────
# A JSON-ban _meta csak gépi adatokat tartalmaz; az útmutató itt él.
CATALOG_GUIDE_TEMPLATE = """\
# figure_catalog.json — Útmutató

**Részletes szabályok:** `.claude/skills/02b_figure_enricher.md`

---

## Prefix-konvenció

```
_ prefix  = script-kezelt mező, ne szerkeszd kézzel
  Példák:  _status (derived), _meta (gépi metaadat)

Nincs _   = user-editable
  Példák:  notes, keywords, caption, caption_verified, visual_content, text_context
```

---

## Szerepkörök

| Jelölés | Ki tölti |
|---------|----------|
| 🐍 | Python script automatikusan |
| 🤖 | Claude (02b skill) |
| 😎 | User verifikálja vagy javítja |
| 🤖→😎 | Claude javasol, te véglegesíted |

---

## `_status` értékek (🐍 derived — NE szerkeszd kézzel)

| Érték | Feltétel | Teendő |
|-------|----------|--------|
| `complete` | `caption_verified:true` ÉS `visual_content` kitöltve | Kész, 05 retrieval használhatja |
| `caption-ok` | `caption_verified:true`, de `visual_content:null` | 02b bootstrap hiányzik |
| `draft` | `visual_content` kitöltve, de `caption_verified:false` | 😎 jóváhagyás hiányzik |
| `un-processed` | Sem verified, sem visual_content | 02b még nem futott |

---

## Értékkonvenciók

| JSON érték | Jelentés |
|------------|---------|
| `null` | Feldolgozás nem futott le (`un-processed`) |
| `[]` | Feldolgozva, de üres eredmény (`no-results`) |
| `true/false` | Boolean: tudottan pozitív/negatív |

**Tilos:** `""` (üres string) — helyette `null`.

---

## Mezők

| Mező | Ki | Leírás |
|------|----|--------|
| `id` | 🐍 | Egyedi azonosító (fig_NNN). Stabil lookup-kulcs: `(source_file, path)` |
| `page` | 🐍 | Forrás-oldal száma |
| `path` | 🐍 | Kép path `2_clean_inputs/<src>/images/pNNN_figNNN.png` alatt |
| `needs_crop` | 🐍→😎 | `true` = még vágni kell; `false` = kész / nem kellett |
| `caption` | 🤖→😎 | Az ábra eredeti felirata, paragrafus-szennyezés nélkül |
| `caption_verified` | 😎 | `true`, ha vizuálisan ÉS szövegileg megerősítetted |
| `visual_content` | 🤖→😎 | 1-3 mondat: mit ábrázol (diagram, tengelyek, fő elemek) |
| `text_context` | 🤖→😎 | 1-3 mondat: mi a szöveg-környezet lényege |
| `keywords` | 🤖→😎 | 3-8 kulcsszó; logók esetén tartalmazzon `"logo"` tag-et |
| `_status` | 🐍 | Derived állapot — lásd fenti táblázat. NE szerkeszd kézzel |
| `notes` | 😎 | Szabad szövegű megjegyzések listája. `"✅ reviewed"` = teljes átnézés |

**Megjegyzések írása:**
- Egy konkrét ábrához → az entry `notes` listájába
- Sprint/folyamat szinten → `.claude/sprints/image_rag/review_notes.md`

---

## Példa — teljesen kitöltött entry

```json
{
  "id": "fig_000_EXAMPLE",
  "page": 7,
  "path": "2_clean_inputs/example_paper/images/p007_fig001.png",
  "needs_crop": false,
  "caption": "Figure 2: Sample compressor map showing surge line and operating point.",
  "caption_verified": true,
  "visual_content": "Kompresszor-jelleggörbe: PR a Y-tengelyen, Mass flow az X-en. Bal felső sarokban a surge line; egy A pont a stabil tartományban.",
  "text_context": "A 3. SURGE FUNDAMENTALS szekció elején, a surge line bevezetésénél. Hivatkozás: p6 (Section 3 intro).",
  "keywords": ["surge line", "compressor map", "operating point", "stability boundary"],
  "_status": "complete",
  "notes": ["Az A pont feliratát meg lehetett volna nagyobbra venni.", "✅ reviewed"]
}
```
"""


# ── Helperek (séma-tudatosak) ──────────────────────────────────────────────────

def load_citations(week_dir: Path) -> dict:
    """citations.json betöltése fájlnév→citáció-kulcs mappinghoz."""
    cit_path = week_dir / "1_raw_inputs" / "citations.json"
    if not cit_path.exists():
        return {}
    data = json.loads(cit_path.read_text(encoding="utf-8"))
    return {v["filename"]: k for k, v in data.items()
            if k != "_meta" and v.get("filename")}


def new_catalog() -> dict:
    """Üres v4 katalógus váz. _meta = gépi adatok; útmutató → CATALOG_GUIDE.md."""
    return {
        "_meta": {
            "schema_version": SCHEMA_VERSION,
            "last_updated": date.today().isoformat(),
            "_guide": "CATALOG_GUIDE.md",
        },
        "sources": {},
    }


def load_catalog(path: Path) -> dict:
    """v4 katalógus betöltése. Más sémát hard error-ral elutasít."""
    if not path.exists():
        return new_catalog()
    raw = json.loads(path.read_text(encoding="utf-8"))
    schema = raw.get("_meta", {}).get("schema_version") if isinstance(raw, dict) else None
    if schema != SCHEMA_VERSION:
        sys.exit(
            f"HIBA: nem v{SCHEMA_VERSION} séma ({schema}) a katalógusban {path}. "
            "Wipe + regen, vagy explicit migráció szükséges."
        )
    # Defenzív: hiányzó mezők kitöltése
    for src_data in raw.get("sources", {}).values():
        for e in src_data.get("figures", []):
            for k, default in ENTRY_DEFAULTS.items():
                if k not in e:
                    e[k] = default.copy() if isinstance(default, list) else default
    return raw


def write_catalog_guide(clean_in: Path, dry_run: bool) -> None:
    """CATALOG_GUIDE.md generálása a katalógus mellé (ha még nem létezik)."""
    guide_path = clean_in / "CATALOG_GUIDE.md"
    if guide_path.exists():
        return  # idempotens: meglevőt nem írjuk felül
    if not dry_run:
        guide_path.write_text(CATALOG_GUIDE_TEMPLATE, encoding="utf-8")
        print(f"  📄 CATALOG_GUIDE.md generálva → {guide_path}")


def _ensure_source(catalog: dict, source_file: str, citation_key: str) -> dict:
    """Visszaadja (létrehozza ha kell) a source_file blokkot."""
    if source_file not in catalog["sources"]:
        catalog["sources"][source_file] = {
            "citation_key": str(citation_key),
            "figures": [],
        }
    return catalog["sources"][source_file]


def all_figures(catalog: dict):
    """Flat iterator minden bejegyzésen (minden forrás)."""
    for src_data in catalog["sources"].values():
        for entry in src_data["figures"]:
            yield entry


def next_fig_id(catalog: dict) -> str:
    """Következő fig_NNN azonosító a katalógus összes bejegyzése alapján."""
    existing = [int(e["id"].split("_")[1])
                for e in all_figures(catalog)
                if e.get("id", "").startswith("fig_") and "_" in e["id"]]
    n = max(existing, default=0) + 1
    return f"fig_{n:03d}"


def already_in_catalog(catalog: dict, source_file: str, path_str: str) -> bool:
    """(source_file, path) pár szerepel-e már a katalógusban?"""
    src_data = catalog["sources"].get(source_file)
    if not src_data:
        return False
    return any(e.get("path") == path_str for e in src_data["figures"])


def make_entry(fig_id: str, page: int, path_str: str, needs_crop: bool) -> dict:
    """Új bejegyzés strukturális mezőkkel + alapérték meta-mezőkkel.
    LOGIKAI sorrendben — JSON-ban ugyanígy fog megjelenni."""
    entry: dict = {}
    for k, default in ENTRY_DEFAULTS.items():
        entry[k] = default.copy() if isinstance(default, (list, dict)) else default
    entry["id"] = fig_id
    entry["page"] = page
    entry["path"] = path_str
    entry["needs_crop"] = needs_crop
    return entry


def _compute_status(entry: dict) -> str:
    """Származtatott 4-állapotú státusz. _ prefix = script-managed, ne szerkeszd.
    Lásd CATALOG_GUIDE.md _status táblázat."""
    caption_ok = bool(entry.get("caption_verified"))
    has_meta   = bool(entry.get("visual_content"))
    if caption_ok and has_meta:   return "complete"
    if caption_ok:                return "caption-ok"
    if has_meta:                  return "draft"
    return "un-processed"


def _refresh_statuses(catalog: dict) -> None:
    """Minden bejegyzésen újraszámolja a _status mezőt."""
    for e in all_figures(catalog):
        e["_status"] = _compute_status(e)


def save_catalog(catalog: dict, path: Path, dry_run: bool):
    if dry_run:
        return
    _refresh_statuses(catalog)
    catalog["_meta"]["last_updated"] = date.today().isoformat()
    path.write_text(json.dumps(catalog, ensure_ascii=False, indent=2),
                    encoding="utf-8")
    write_catalog_guide(path.parent, dry_run)


# ── OCR helper ─────────────────────────────────────────────────────────────────

def _try_ocr_page(page, pixmap, out_dir: Path, page_num: int,
                  dry_run: bool) -> Path | None:
    """OCR-ezi a szkennelt oldalt, ha szükséges és lehetséges.
    Logic:
      1. Cache hit: ha a .txt már létezik → return path (idempotens)
      2. Born-digital ellenőrzés: page.get_text() nem üres → no OCR needed
      3. Tesseract elérhetőség: ha nincs → WARN, return None
      4. OCR + mentés text/pNNN.txt-be
    """
    txt_path = out_dir / "text" / f"p{page_num:03d}.txt"
    if txt_path.exists() and txt_path.stat().st_size > 0:
        return txt_path
    # Born-digital text stream check
    try:
        if page.get_text("text").strip():
            return None  # van text stream, OCR felesleges
    except Exception:
        pass
    if not _OCR_AVAILABLE:
        print(f"  ⚠️  OCR kihagyva (pytesseract / Pillow / Tesseract nem elérhető): "
              f"oldal {page_num}")
        return None
    try:
        from io import BytesIO
        img = Image.open(BytesIO(pixmap.tobytes("png")))
        text = pytesseract.image_to_string(img, lang=OCR_LANGS)
    except Exception as e:
        print(f"  ⚠️  OCR hiba oldal {page_num}: {e}", file=sys.stderr)
        return None
    if not dry_run:
        txt_path.parent.mkdir(parents=True, exist_ok=True)
        txt_path.write_text(text, encoding="utf-8")
    print(f"  📝 OCR: oldal {page_num} → {txt_path.relative_to(out_dir.parent.parent.parent)} "
          f"({len(text)} char)")
    return txt_path


# ── PDF extractor ──────────────────────────────────────────────────────────────

def _img_name(page_num: int, fig_idx: int) -> str:
    """Egységes fájlnév-konvenció: pNNN_figNNN.png mindenhol."""
    return f"p{page_num:03d}_fig{fig_idx:03d}.png"


def _rel_path(src: Path, img_name: str) -> str:
    return f"2_clean_inputs/{src.stem}/images/{img_name}"


def extract_pdf(src: Path, out_dir: Path, citation_key: str,
                catalog: dict, dry_run: bool,
                specific_pages: list[int] | None = None) -> tuple[int, int, int]:
    """Born-digital + vegyes + szkennelt PDF kezelés.
    Visszatérés: (mentett, kihagyott_deko, crop_figyelmeztetések)
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print(f"  HIBA: PyMuPDF nincs telepítve. pip install pymupdf", file=sys.stderr)
        return 0, 0, 0

    doc = fitz.open(str(src))
    n_pages = len(doc)

    # ── specific_pages mód: Claude-azonosított oldalak ─────────────────────────
    if specific_pages is not None:
        page_counts = Counter(specific_pages)
        saved = skipped = crop_warn = 0

        for page_num, fig_count in sorted(page_counts.items()):
            if page_num < 1 or page_num > n_pages:
                print(f"  SKIP  oldal {page_num}: kívül esik ({n_pages} old.)")
                continue
            page = doc[page_num - 1]
            page_area = page.rect.width * page.rect.height
            page_imgs = page.get_images(full=True)

            is_scanned = any(
                (doc.extract_image(img[0])["width"] * doc.extract_image(img[0])["height"])
                / page_area >= PAGE_FILL
                for img in page_imgs[:1]
            ) if page_imgs else False

            if is_scanned:
                # Szkennelt oldal: 1 render → N catalog entry (mind ugyanarra)
                mat = fitz.Matrix(RENDER_DPI / 72, RENDER_DPI / 72)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                # OCR cache (egy oldalrender → egy txt)
                _try_ocr_page(page, pix, out_dir, page_num, dry_run)
                for fig_idx in range(1, fig_count + 1):
                    img_name = _img_name(page_num, fig_idx)
                    rel = _rel_path(src, img_name)
                    print(f"  ⚠️  CROP SZÜKSÉGES: oldal {page_num} kép {fig_idx} → {img_name}")
                    if not dry_run:
                        img_path = out_dir / "images" / img_name
                        img_path.parent.mkdir(parents=True, exist_ok=True)
                        img_path.write_bytes(img_bytes)
                    if not already_in_catalog(catalog, src.name, rel):
                        _ensure_source(catalog, src.name, citation_key)["figures"].append(
                            make_entry(next_fig_id(catalog), page_num, rel,
                                       needs_crop=True))
                        saved += 1
                        crop_warn += 1
                    else:
                        skipped += 1
            else:
                # Born-digital oldal: egységes pNNN_figNNN.png, folyamatos counter
                page_fig_idx = 0
                # 1. Beágyazott raszterek
                for img_info in page_imgs:
                    xref = img_info[0]
                    base = doc.extract_image(xref)
                    w, h = base["width"], base["height"]
                    if w * h < MIN_AREA:
                        skipped += 1
                        continue
                    page_fig_idx += 1
                    img_name = _img_name(page_num, page_fig_idx)
                    rel = _rel_path(src, img_name)
                    img_bytes = base["image"]
                    if base["ext"] != "png":
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n > 4:
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                        img_bytes = pix.tobytes("png")
                    if not dry_run:
                        img_path = out_dir / "images" / img_name
                        img_path.parent.mkdir(parents=True, exist_ok=True)
                        img_path.write_bytes(img_bytes)
                    if not already_in_catalog(catalog, src.name, rel):
                        _ensure_source(catalog, src.name, citation_key)["figures"].append(
                            make_entry(next_fig_id(catalog), page_num, rel,
                                       needs_crop=False))
                        saved += 1
                    else:
                        skipped += 1

                # 2. Ha nincs raszter, vektoros oldal-render
                has_raster = page_fig_idx > 0
                if not has_raster:
                    mat = fitz.Matrix(RENDER_DPI / 72, RENDER_DPI / 72)
                    pix = page.get_pixmap(matrix=mat)
                    img_bytes = pix.tobytes("png")
                    for fig_idx in range(1, fig_count + 1):
                        page_fig_idx += 1
                        img_name = _img_name(page_num, page_fig_idx)
                        rel = _rel_path(src, img_name)
                        needs_crop = True
                        if not dry_run:
                            img_path = out_dir / "images" / img_name
                            img_path.parent.mkdir(parents=True, exist_ok=True)
                            img_path.write_bytes(img_bytes)
                            if _auto_crop_fn is not None:
                                cropped, ratio = _auto_crop_fn(img_path)
                                if cropped:
                                    needs_crop = False
                                    print(f"  ✂️  AUTO-CROP: oldal {page_num} kép {fig_idx} → {img_name} ({ratio:.0%} levágva)")
                                else:
                                    print(f"  ⚠️  CROP SZÜKSÉGES (vektor): oldal {page_num} kép {fig_idx} → {img_name}")
                            else:
                                print(f"  ⚠️  CROP SZÜKSÉGES (vektor): oldal {page_num} kép {fig_idx} → {img_name}")
                        else:
                            print(f"  ⚠️  CROP SZÜKSÉGES (vektor): oldal {page_num} kép {fig_idx} → {img_name}")
                        if not already_in_catalog(catalog, src.name, rel):
                            _ensure_source(catalog, src.name, citation_key)["figures"].append(
                                make_entry(next_fig_id(catalog), page_num, rel,
                                           needs_crop=needs_crop))
                            saved += 1
                            if needs_crop:
                                crop_warn += 1
                        else:
                            if not needs_crop:
                                for e in all_figures(catalog):
                                    if e.get("path") == rel:
                                        e["needs_crop"] = False
                                        break
                            skipped += 1
        doc.close()
        return saved, skipped, crop_warn

    # ── Normál mód: végigmegy az összes oldalon ────────────────────────────────
    # 1. átmenet: szkennelt oldalak detektálása
    scanned_pages = []
    for page_num, page in enumerate(doc, 1):
        page_area = page.rect.width * page.rect.height
        for img_info in page.get_images(full=True):
            base = doc.extract_image(img_info[0])
            if (base["width"] * base["height"]) / page_area >= PAGE_FILL:
                scanned_pages.append(page_num)
                break

    scanned_ratio = len(scanned_pages) / n_pages if n_pages else 0
    if scanned_ratio >= SCANNED_THRESHOLD:
        print(f"  ⚠️  SZKENNELT FORRÁS ({len(scanned_pages)}/{n_pages} oldal, "
              f"{scanned_ratio:.0%}) — ábra-kinyerés kihagyva. "
              f"Használd: --source {src.name} --pages <oldalszámok> (Claude azonosítja)")
        doc.close()
        return 0, 0, 0

    # 2. átmenet: tényleges feldolgozás
    saved = skipped = crop_warn = 0
    for page_num, page in enumerate(doc, 1):
        page_area = page.rect.width * page.rect.height
        is_scanned_page = page_num in scanned_pages
        page_fig_idx = 0
        page_imgs = page.get_images(full=True)

        # Vegyes-PDF szkennelt oldal: 1 render → 1 catalog entry
        # (függetlenül attól, hány embedded image van a page_imgs-ben)
        if is_scanned_page:
            mat = fitz.Matrix(RENDER_DPI / 72, RENDER_DPI / 72)
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")
            _try_ocr_page(page, pix, out_dir, page_num, dry_run)
            page_fig_idx = 1
            img_name = _img_name(page_num, page_fig_idx)
            rel = _rel_path(src, img_name)
            print(f"  ⚠️  CROP SZÜKSÉGES (oldal): oldal {page_num} → {img_name}")
            if not dry_run:
                img_path = out_dir / "images" / img_name
                img_path.parent.mkdir(parents=True, exist_ok=True)
                img_path.write_bytes(img_bytes)
            if not already_in_catalog(catalog, src.name, rel):
                _ensure_source(catalog, src.name, citation_key)["figures"].append(
                    make_entry(next_fig_id(catalog), page_num, rel, needs_crop=True))
                saved += 1
                crop_warn += 1
            else:
                skipped += 1
            continue  # ne dolgozzuk fel a embedded raszter-listát szkennelt oldalon

        # Born-digital oldal: végigmegy az embedded raszterek listáján
        for img_info in page_imgs:
            xref = img_info[0]
            base = doc.extract_image(xref)
            w, h = base["width"], base["height"]
            area = w * h
            if area < MIN_AREA:
                skipped += 1
                continue
            page_fig_idx += 1
            img_name = _img_name(page_num, page_fig_idx)
            img_bytes = base["image"]
            if base["ext"] != "png":
                pix = fitz.Pixmap(doc, xref)
                if pix.n > 4:
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                img_bytes = pix.tobytes("png")

            rel = _rel_path(src, img_name)
            if not dry_run:
                img_path = out_dir / "images" / img_name
                img_path.parent.mkdir(parents=True, exist_ok=True)
                img_path.write_bytes(img_bytes)

            if not already_in_catalog(catalog, src.name, rel):
                _ensure_source(catalog, src.name, citation_key)["figures"].append(
                    make_entry(next_fig_id(catalog), page_num, rel, needs_crop=False))
            saved += 1

        # Vektoros oldal-render (ha sem raszter, sem szkennelt)
        if not is_scanned_page and page_fig_idx == 0:
            page_w = page.rect.width
            sig_drawings = [
                d for d in page.get_drawings()
                if d["rect"].width * d["rect"].height >= VECTOR_MIN_ELEM_AREA
                and not (d["rect"].height < VECTOR_DIVIDER_H
                         and d["rect"].width > 0.6 * page_w)
            ]
            if len(sig_drawings) >= VECTOR_MIN_DRAW:
                mat = fitz.Matrix(RENDER_DPI / 72, RENDER_DPI / 72)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                page_fig_idx += 1
                img_name = _img_name(page_num, page_fig_idx)
                rel = _rel_path(src, img_name)
                needs_crop = True
                if _auto_crop_fn is not None and not dry_run:
                    img_path = out_dir / "images" / img_name
                    img_path.parent.mkdir(parents=True, exist_ok=True)
                    img_path.write_bytes(img_bytes)
                    cropped, ratio = _auto_crop_fn(img_path)
                    if cropped:
                        needs_crop = False
                        print(f"  ✂️  AUTO-CROP (vektor): oldal {page_num} → {img_name} ({ratio:.0%} levágva)")
                    else:
                        print(f"  ⚠️  CROP SZÜKSÉGES (vektor): oldal {page_num} → {img_name} ({len(sig_drawings)} drawing)")
                        crop_warn += 1
                else:
                    print(f"  ⚠️  CROP SZÜKSÉGES (vektor): oldal {page_num} → {img_name} ({len(sig_drawings)} drawing)")
                    if not dry_run:
                        img_path = out_dir / "images" / img_name
                        img_path.parent.mkdir(parents=True, exist_ok=True)
                        img_path.write_bytes(img_bytes)
                    crop_warn += 1
                if not already_in_catalog(catalog, src.name, rel):
                    _ensure_source(catalog, src.name, citation_key)["figures"].append(
                        make_entry(next_fig_id(catalog), page_num, rel, needs_crop=needs_crop))
                    saved += 1

    doc.close()
    return saved, skipped, crop_warn


# ── PPTX extractor ─────────────────────────────────────────────────────────────

def _collect_pptx_images(shapes) -> list:
    """Rekurzív gyűjtés minden blipFill-tartalmú shape-ből (group-on belül is)."""
    result = []
    for shape in shapes:
        if shape.shape_type == 6 and hasattr(shape, "shapes"):
            result.extend(_collect_pptx_images(shape.shapes))
            continue
        try:
            from lxml import etree
            xml = etree.tostring(shape.element).decode("utf-8")
        except Exception:
            xml = getattr(shape.element, "xml", "")
        if "blipFill" not in xml and "a:blip" not in xml:
            continue
        result.append(shape)
    return result


def extract_pptx(src: Path, out_dir: Path, citation_key: str,
                 catalog: dict, dry_run: bool) -> tuple[int, int]:
    """PPTX → pNNN_figNNN.png (slide N = page N)."""
    try:
        from pptx import Presentation
    except ImportError:
        print(f"  HIBA: python-pptx nincs telepítve. pip install python-pptx", file=sys.stderr)
        return 0, 0

    prs = Presentation(str(src))
    saved = skipped = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        img_shapes = _collect_pptx_images(slide.shapes)
        page_fig_idx = 0
        for shape in img_shapes:
            try:
                image = shape.image
                blob = image.blob
                w = shape.width.pt if hasattr(shape.width, "pt") else 0
                h = shape.height.pt if hasattr(shape.height, "pt") else 0
                area = w * h
                if area < MIN_AREA / 10:
                    skipped += 1
                    continue
                page_fig_idx += 1
                img_name = _img_name(slide_idx, page_fig_idx)
                rel = _rel_path(src, img_name)

                # PNG-vé konvertálás
                try:
                    import fitz
                    pix = fitz.Pixmap(blob)
                    if pix.n > 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    img_bytes = pix.tobytes("png")
                except Exception:
                    img_bytes = blob

                if not dry_run:
                    img_path = out_dir / "images" / img_name
                    img_path.parent.mkdir(parents=True, exist_ok=True)
                    img_path.write_bytes(img_bytes)

                if not already_in_catalog(catalog, src.name, rel):
                    _ensure_source(catalog, src.name, citation_key)["figures"].append(
                        make_entry(next_fig_id(catalog), slide_idx, rel, needs_crop=False))
                    saved += 1
            except Exception as e:
                print(f"  SKIP  slide {slide_idx} kep: {e}", file=sys.stderr)

    return saved, skipped


# ── Főfüggvény ─────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Ábra-kinyerő (v4 séma): PDF/PPTX → pNNN_figNNN.png + figure_catalog.json")
    parser.add_argument("--week-dir", required=True, type=Path,
                        help="Heti mappa (pl. test_outputs/atg/1_het)")
    parser.add_argument("--source", type=str, default=None,
                        help="Forrás fájlneve (csak ehhez futtat, --pages-szel)")
    parser.add_argument("--pages", type=str, default=None,
                        help="Claude-azonosított oldalszámok vesszővel (pl. '5,12,23')")
    parser.add_argument("--dry-run", action="store_true",
                        help="Listázás mentés nélkül")
    args = parser.parse_args()

    week_dir = args.week_dir.resolve()
    raw_in   = week_dir / "1_raw_inputs"
    clean_in = week_dir / "2_clean_inputs"
    cat_path = clean_in / "figure_catalog.json"

    if not raw_in.is_dir():
        sys.exit(f"HIBA: nem található {raw_in}")

    catalog = load_catalog(cat_path)
    citations = load_citations(week_dir)
    prefix = "[DRY] " if args.dry_run else ""

    specific_pages = None
    if args.source and args.pages:
        try:
            specific_pages = [int(p.strip()) for p in args.pages.split(",") if p.strip()]
        except ValueError:
            sys.exit("HIBA: --pages csak egész számokat fogad el, vesszővel (pl. '5,12,23')")

    total_saved = total_skip = total_crop = 0
    processed = []

    if args.source:
        src_path = raw_in / args.source
        if not src_path.exists():
            sys.exit(f"HIBA: nem található {src_path}")
        sources = [src_path]
    else:
        sources = sorted(f for f in raw_in.iterdir()
                         if f.is_file()
                         and f.suffix.lower() in (".pdf", ".pptx")
                         and not f.name.startswith("_")
                         and f.name != "citations.json")

    for src in sources:
        ext = src.suffix.lower()
        citation_key = citations.get(src.name, "?")
        out_dir = clean_in / src.stem
        label = f"oldal {specific_pages}" if specific_pages else ""
        print(f"{prefix}{src.name} [cit:{citation_key}]{' ' + label if label else ''}")

        if ext == ".pdf":
            saved, skip, crop = extract_pdf(src, out_dir, citation_key,
                                            catalog, args.dry_run, specific_pages)
            total_crop += crop
        elif ext == ".pptx":
            saved, skip = extract_pptx(src, out_dir, citation_key, catalog, args.dry_run)
            crop = 0
        else:
            continue

        total_saved += saved
        total_skip  += skip
        processed.append(src.name)
        print(f"  → {saved} kép mentve, {skip} dekoráció kihagyva"
              + (f", {crop} oldal crop-ra vár ⚠️" if crop else ""))

    save_catalog(catalog, cat_path, args.dry_run)

    print(f"\n{prefix}Kész: {len(processed)} forrás | "
          f"{total_saved} kép | {total_skip} deko kihagyva"
          + (f" | ⚠️  {total_crop} oldal CROP SZÜKSÉGES" if total_crop else ""))

    if total_crop:
        print("  Crop-ra váró bejegyzések a figure_catalog.json-ban: needs_crop: true")
        pending = [e["id"] for e in all_figures(catalog) if e.get("needs_crop")]
        if pending:
            print(f"  Pending crop IDs: {', '.join(pending)}")


if __name__ == "__main__":
    main()
