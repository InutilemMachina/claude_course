"""
_crop_tasks.py — Crop-feladatlista generálás és checkbox→catalog szinkron.

Standalone utility; a 02_image_extraction.py dynamikusan tölti be.

Funkciók:
    generate_crop_tasks(week_dir)  — figure_catalog.json → _crop_tasks.md
                                    (forrásonként markdown táblázat,
                                     minden bejegyzés, [x]/[ ] + Caption oszlop)
    sync_crop_tasks(week_dir)      — [x] és Caption cellákat catalog-ba menti,
                                    majd regenerálja a _crop_tasks.md-t.

Caption auto-detekció:
    - Born-digital PDF: page text → `Figure N(.M)?: ...` minta.
    - PPTX: slide text → ugyanaz a minta.
    - Szkennelt PDF / üres: caption üres.
    - Bizonytalan találat (több caption az oldalon, index-mismatch): a Caption
      cella elején `?` jel — a megerősítéshez töröld a `?`-et.
"""

import json
import re
from datetime import date
from pathlib import Path


# ── Belső segédek ──────────────────────────────────────────────────────────────

def _load_catalog(clean_in: Path) -> list:
    cat_path = clean_in / "figure_catalog.json"
    if not cat_path.exists():
        raise FileNotFoundError(f"Nem található: {cat_path}")
    return json.loads(cat_path.read_text(encoding="utf-8"))


def _save_catalog(catalog: list, clean_in: Path) -> None:
    cat_path = clean_in / "figure_catalog.json"
    cat_path.write_text(json.dumps(catalog, ensure_ascii=False, indent=2),
                        encoding="utf-8")


def _cit_key(source_file: str, catalog: list) -> str:
    """A forráshoz tartozó citation_key (az első bejegyzésből)."""
    for e in catalog:
        if e.get("source_file") == source_file:
            return str(e.get("citation_key", "?"))
    return "?"


def _week_label(week_dir: Path) -> str:
    """'atg/1_het' stílusú label a week_dir abszolút útjából."""
    parts = week_dir.parts
    if len(parts) >= 2:
        return f"{parts[-2]}/{parts[-1]}"
    return week_dir.name


# ── Caption auto-detekció ──────────────────────────────────────────────────────

_CAPTION_RE = re.compile(
    r"(?:Figure|Fig\.?)\s+(\d+(?:\.\d+)?)\s*[:\.\)]\s*([^.\n]{3,250}\.)",
    re.IGNORECASE,
)

_IDX_RE = re.compile(r"(?:fig|img)(\d+)", re.IGNORECASE)


def _pdf_captions(pdf_path: Path) -> dict:
    """{page_number: [caption strings in document order]}"""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return {}
    try:
        doc = fitz.open(pdf_path)
    except Exception:
        return {}
    cache: dict = {}
    for i, page in enumerate(doc, start=1):
        try:
            text = page.get_text("text") or ""
        except Exception:
            text = ""
        flat = re.sub(r"\s+", " ", text).strip()
        matches = _CAPTION_RE.findall(flat)
        cache[i] = [f"Figure {num}: {body.strip()}" for num, body in matches]
    doc.close()
    return cache


def _pptx_captions(pptx_path: Path) -> dict:
    """{slide_number: [caption strings]}"""
    try:
        from pptx import Presentation
    except ImportError:
        return {}
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return {}
    cache: dict = {}
    for i, slide in enumerate(prs.slides, start=1):
        chunks = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for p in shape.text_frame.paragraphs:
                    chunks.append(p.text or "")
        flat = re.sub(r"\s+", " ", " ".join(chunks)).strip()
        matches = _CAPTION_RE.findall(flat)
        cache[i] = [f"Figure {num}: {body.strip()}" for num, body in matches]
    return cache


def _build_caption_cache(week_dir: Path, sources: list) -> dict:
    """source_file → {page: [captions]}"""
    raw_dir = week_dir / "1_raw_inputs"
    cache: dict = {}
    for src in sources:
        path = raw_dir / src
        if not path.exists():
            cache[src] = {}
            continue
        ext = path.suffix.lower()
        if ext == ".pdf":
            cache[src] = _pdf_captions(path)
        elif ext == ".pptx":
            cache[src] = _pptx_captions(path)
        else:
            cache[src] = {}
    return cache


def _caption_for(entry: dict, cap_cache: dict) -> tuple:
    """
    Visszatérés: (caption_text, uncertain).
    1. Ha catalog.caption ki van töltve → azt használjuk (uncertain=False).
    2. Egyébként a forrás-cache-ből keresünk az oldal + image_index alapján.
    3. Bizonytalanság: több caption az oldalon, vagy index out-of-range.
    """
    cat_cap = entry.get("caption")
    if cat_cap:
        return str(cat_cap), False
    src = entry.get("source_file", "")
    page = entry.get("page")
    page_caps = cap_cache.get(src, {}).get(page, [])
    if not page_caps:
        return "", False
    fname = Path(entry.get("filename", "")).name
    m = _IDX_RE.search(fname)
    idx = int(m.group(1)) if m else 1
    if 1 <= idx <= len(page_caps):
        uncertain = len(page_caps) > 1  # 1 image / 1 caption → confident
        return page_caps[idx - 1], uncertain
    return page_caps[0], True


# ── Publikus API ───────────────────────────────────────────────────────────────

def _escape_cell(text: str) -> str:
    """Markdown table cell escaping (pipe + newline)."""
    return text.replace("\\", "\\\\").replace("|", "\\|").replace("\n", " ").strip()


def generate_crop_tasks(week_dir: Path) -> None:
    """
    figure_catalog.json → _crop_tasks.md.

    Forrásonként egy markdown táblázat, MINDEN bejegyzéssel:
        | ✓ | id | fájl | oldal | útvonal | Caption |
    `[x]` ha `needs_crop: false`, `[ ]` ha `true`.
    Caption cella üres ha nincs találat, `? ` prefixszel ha auto-detektált,
    de bizonytalan (több caption az oldalon vagy index-mismatch).
    """
    week_dir = Path(week_dir).resolve()
    clean_in = week_dir / "2_clean_inputs"
    md_path  = clean_in / "_crop_tasks.md"

    catalog = _load_catalog(clean_in)
    if not catalog:
        md_path.write_text("# Crop tasks\n\nÜres katalógus.\n",
                           encoding="utf-8")
        print("  _crop_tasks.md: üres katalógus")
        return

    sources_ordered = sorted({e["source_file"] for e in catalog})
    cap_cache = _build_caption_cache(week_dir, sources_ordered)

    pending = sum(1 for e in catalog if e.get("needs_crop"))
    total   = len(catalog)
    today   = date.today().isoformat()
    label   = _week_label(week_dir)

    lines = []
    lines.append(f"# Crop tasks — {label}")
    lines.append(f"_{pending} crop vár / {total} összesen | Frissítve: {today}_")
    lines.append("")
    lines.append("> **Munkamenet:** `[ ]` = még crop-olni kell; `[x]` = kész.")
    lines.append("> A `Caption` oszlopba írd a felirat szövegét. `?` prefix = a script auto-detektálta, de bizonytalan — ellenőrizd és töröld a `?`-et a megerősítéshez.")
    lines.append("> Üres Caption cella = nem találtunk feliratot, töltsd ki kézzel.")
    lines.append("")
    lines.append("---")

    for src in sources_ordered:
        entries = [e for e in catalog if e["source_file"] == src]
        cit_k   = _cit_key(src, catalog)
        lines.append("")
        lines.append(f"## {src}  [cit:{cit_k}]")
        lines.append("")
        lines.append("| ✓ | id | fájl | oldal | útvonal | Caption |")
        lines.append("|:-:|:--|:--|:-:|:--|:--|")
        for e in entries:
            fig_id   = e["id"]
            filename = Path(e["filename"]).name
            rel_path = e["filename"]
            page     = e.get("page", "?")
            checkbox = "[x]" if not e.get("needs_crop") else "[ ]"
            cap_text, uncertain = _caption_for(e, cap_cache)
            cap_cell = _escape_cell(cap_text)
            if cap_cell and uncertain:
                cap_cell = f"? {cap_cell}"
            lines.append(
                f"| {checkbox} | {fig_id} | `{filename}` | {page} "
                f"| `{rel_path}` | {cap_cell} |"
            )

    lines.append("")   # záró newline

    md_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"  _crop_tasks.md: {total} bejegyzés ({pending} crop vár) → {md_path}")


# Markdown-táblázat sor-parser a sync-hez
_CHECKBOX_RE = re.compile(r"\[\s*([ xX])\s*\]")
_FIGID_RE    = re.compile(r"^fig_\w+$")


def _split_row(line: str) -> list:
    """Markdown táblázat sort cellákra bont (escape-elt `\\|` megőrzve)."""
    stripped = line.strip()
    if not stripped.startswith("|"):
        return []
    # `\|` ideiglenes placeholderre cseréljük, hogy a split ne tönkretegye
    PLACEHOLDER = "\x00PIPE\x00"
    s = stripped.strip("|").replace("\\|", PLACEHOLDER)
    return [c.replace(PLACEHOLDER, "|").strip() for c in s.split("|")]


def sync_crop_tasks(week_dir: Path) -> int:
    """
    _crop_tasks.md → figure_catalog.json:
        - `[x]` checkbox → `needs_crop: false`.
        - Caption cella (NEM `?`-prefixes) → `caption: <text>` a catalog-ban.

    Bizonytalan (`?` prefix) caption-eket NEM mentjük — azok a következő
    generálásnál újra auto-detektálódnak, amíg a user a `?`-et el nem távolítja.

    Végül regenerálja a _crop_tasks.md-t (a kész [x]-ek bennmaradnak).
    Visszatérés: módosított bejegyzések száma.
    """
    week_dir = Path(week_dir).resolve()
    clean_in = week_dir / "2_clean_inputs"
    md_path  = clean_in / "_crop_tasks.md"

    if not md_path.exists():
        print(f"  WARN: nem található: {md_path} — sync kihagyva", flush=True)
        return 0

    text = md_path.read_text(encoding="utf-8")
    catalog = _load_catalog(clean_in)
    by_id = {e["id"]: e for e in catalog}

    updated_crop = 0
    updated_cap  = 0

    for raw_line in text.splitlines():
        cells = _split_row(raw_line)
        if len(cells) < 6:
            continue
        cb_cell, fig_id, _fname, _page, _path, cap_cell = cells[:6]

        cb_m = _CHECKBOX_RE.search(cb_cell)
        if not cb_m:
            continue
        if not _FIGID_RE.match(fig_id):
            continue   # header vagy separator sor

        entry = by_id.get(fig_id)
        if not entry:
            continue

        checked = cb_m.group(1).strip().lower() == "x"
        if checked and entry.get("needs_crop"):
            entry["needs_crop"] = False
            updated_crop += 1

        uncertain = cap_cell.startswith("?")
        cap_text  = cap_cell.lstrip("?").strip()
        if cap_text and not uncertain and entry.get("caption") != cap_text:
            entry["caption"] = cap_text
            updated_cap += 1

    if updated_crop or updated_cap:
        _save_catalog(catalog, clean_in)
        print(f"  sync: needs_crop +{updated_crop}, caption +{updated_cap} "
              f"bejegyzés frissítve")
    else:
        print("  sync: nincs változás")

    generate_crop_tasks(week_dir)   # regenerálás (most már minden bent marad)
    return updated_crop + updated_cap
