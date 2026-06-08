"""
10_pptx_gyarto.py — MARP Markdown → PPTX (két variáns: default / mindmap)

A .potx layoutjait használja — nincs kézi stílusdefiníció a scriptben.
Stílus (font, szín, chrome, logo, footer) a .potx-ből örökl; a script csak a
layout-kiválasztást és a placeholder-feltöltést végzi. A potx↔python szerződés
kizárólag placeholder-idx alapú: idx0=cím, idx1=body, idx2=kép, idx3=felirat,
idx5=mindmap_body.

Variánsok (azonos navigációs modellből, lásd _nav_util.py):
  - default : a tájékozódást a (többsoros) fejléc-breadcrumb adja; nincs oldalsáv.
              Sablon: due_presentation_default_master.potx
  - mindmap : a tájékozódást a jobb oldali sorszámozott TOC (idx5) adja.
              Sablon: due_presentation_mindmap_master.potx

A navigáció mindig SZÖVEG. A meglévő _prezi_assets/(navigator|secN).png képek =
navigációs helyek → a PPTX-ben elhagyjuk (helyettük idx5 TOC / idx0 breadcrumb).
Minden más kép tartalmi ábra → marad. A jegyzetből vett valódi diagramok PNG-k.

Usage:
    python scripts/10_pptx_gyarto.py --week-dir test_outputs/<tárgy>/N_het [--variant both]
    python scripts/10_pptx_gyarto.py <md> [--output <pptx>] [--variant mindmap]
"""

import argparse
import io
import os
import re
import sys
import zipfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    sys.exit("python-pptx not installed.  pip install python-pptx")

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import _nav_util as nav  # noqa: E402
import _omml  # noqa: E402  — natív OMML egyenletek (inline + block)

# ── Sablonok variánsonként ─────────────────────────────────────────────────────

POTX_BY_VARIANT = {
    "default": Path("templates/due_presentation_default_master.potx"),
    "mindmap": Path("templates/due_presentation_mindmap_master.potx"),
}
DEFAULT_POTX = POTX_BY_VARIANT["default"]

# Logikai szerep → layout-név, variánsonként (a MM/Mindmap névkülönbség elrejtve).
LAYOUTS = {
    "default": {
        "COVER": "DUE Cím", "SECTION": "DUE Szakaszfejléc", "TOC": "DUE Tartalom (TOC)",
        "H1": "DUE H1 Fejezet", "H2": "DUE H2 Szakasz", "H3": "DUE H3 Alszakasz",
        "KEP": "DUE Kép+Szöveg", "ABRA": "DUE Ábra", "TABLA": "DUE Táblázat",
        "IROD": "DUE Irodalomjegyzék", "URES": "DUE Üres tartalom",
        "VALTOZAS": "DUE Változásjegyzék",
    },
    "mindmap": {
        "COVER": "DUE Cím", "SECTION": "DUE Mindmap Szakaszfejléc", "TOC": "DUE MM Tartalom (TOC)",
        "H1": "DUE MM H1 Fejezet", "H2": "DUE MM H2 Szakasz", "H3": "DUE MM H3 Alszakasz",
        "KEP": "DUE Kép+Szöveg", "ABRA": "DUE Ábra", "TABLA": "DUE Táblázat",
        "IROD": "DUE MM Irodalomjegyzék", "URES": "DUE MM Üres tartalom",
        "VALTOZAS": "DUE Változásjegyzék",
    },
}


# ── .potx betöltés ────────────────────────────────────────────────────────────

def load_potx(path: str | Path) -> Presentation:
    """Megnyitja a .potx-t Presentation()-ként, a content-type in-memory patchelésével.
    python-pptx nem fogadja el a template content-type-ot, ezért szükséges."""
    CT_POTX = b"presentationml.template.main+xml"
    CT_PPTX = b"presentationml.presentation.main+xml"
    data = Path(path).read_bytes()
    buf = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(data)) as zin, \
         zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            c = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                c = c.replace(CT_POTX, CT_PPTX)
            zout.writestr(item, c)
    buf.seek(0)
    return Presentation(buf)


# ── Layout lookup ─────────────────────────────────────────────────────────────

def get_layout(prs: Presentation, name: str, fallback: str | None = None):
    for ly in prs.slide_layouts:
        if ly.name == name:
            return ly
    if fallback:
        print(f"  ⚠️  Layout nem található: {name!r} — fallback: {fallback!r}")
        return get_layout(prs, fallback)
    print(f"  ⚠️  Layout nem található: {name!r} — első layout")
    return prs.slide_layouts[0]


def role_layout(prs: Presentation, variant: str, role: str):
    """Szerep → layout a variáns táblájából; fallback a variáns saját H2-jére
    (így a mindmap variáns soha nem ejti el a sávot egy sima H2-re)."""
    table = LAYOUTS[variant]
    return get_layout(prs, table.get(role, table["H2"]), fallback=table["H2"])


def slide_phs(slide) -> dict:
    return {ph.placeholder_format.idx: ph for ph in slide.placeholders}


# ── Szöveg-tisztítás ──────────────────────────────────────────────────────────

def clean(text: str) -> str:
    """HTML, markdown, LaTeX, kép-szintaxis, blockquote eltávolítása."""
    text = re.sub(r"<!--.*?-->", "", text, flags=re.DOTALL)
    text = re.sub(r"<[^>]+>", "", text)
    text = re.sub(r"\$\$(.+?)\$\$", r"[\1]", text, flags=re.DOTALL)
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"^\|\s*[-: |]+\|\s*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"^>\s*", "", text, flags=re.MULTILINE)
    text = re.sub(r"^!\[[^\]]*\]\([^)]+\)", "", text, flags=re.MULTILINE)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# ── Placeholder feltöltés ─────────────────────────────────────────────────────

def set_tf(ph, text: str):
    """Body placeholder feltöltése soronként; indentáció → level 0/1/2
    (a .potx bullet-stílusai érvényesülnek)."""
    tf = ph.text_frame
    tf.clear()
    lines = [l for l in text.split("\n") if l.strip()]
    if not lines:
        return
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = re.sub(r"^\s*[-*•]\s+", "", line).strip()   # md lista-jelölő le
        p.level = _level(line)


def _level(line: str) -> int:
    indent = len(line) - len(line.lstrip())
    return 2 if indent >= 4 else (1 if indent >= 2 else 0)


def set_title(ph, text: str):
    """Cím placeholder — többsoros (breadcrumb) támogatással, szint nélkül."""
    tf = ph.text_frame
    tf.clear()
    lines = text.split("\n") if text else [""]
    tf.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line


def insert_img_fit(slide, ph, img_path: str, md_dir: str | None):
    """Kép FIT módban a placeholder területére — arány megtartva, levágás nélkül.
    A képet szabad shape-ként illeszti a placeholder koordinátái közé, középre igazítva.
    ph.insert_picture() SZÁNDÉKOSAN KERÜLENDŐ: az fill/crop módban vágja a képet."""
    path = Path(img_path)
    if not path.is_absolute() and md_dir:
        path = Path(md_dir) / path
    if not path.exists():
        print(f"  ⚠️  Kép nem található: {path}")
        return
    try:
        ph_l, ph_t, ph_w, ph_h = ph.left, ph.top, ph.width, ph.height
        try:
            from PIL import Image as _PIL
            with _PIL.open(str(path)) as _img:
                img_w, img_h = _img.size
        except Exception:
            # Pillow nem elérhető: közvetlen fill-be esünk vissza
            ph.insert_picture(str(path))
            return
        # Scale to fit (letterbox — nincs vágás)
        scale = min(ph_w / img_w, ph_h / img_h)
        new_w = int(img_w * scale)
        new_h = int(img_h * scale)
        # Középre a placeholder területén belül
        left = ph_l + (ph_w - new_w) // 2
        top  = ph_t + (ph_h - new_h) // 2
        slide.shapes.add_picture(str(path), left, top, new_w, new_h)
    except Exception as e:
        print(f"  ⚠️  Kép hiba ({path.name}): {e}")


def set_ph_text(ph, text: str):
    try:
        ph.text = text
    except Exception:
        try:
            ph.text_frame.paragraphs[0].text = text
        except Exception:
            pass


# ── MARP parse ────────────────────────────────────────────────────────────────

def parse_marp(md_text: str) -> list[dict]:
    """MARP Markdown → slide dict lista: {title, subtitle, body, is_cover}."""
    md_text = re.sub(r"^---\n.*?---\n", "", md_text, count=1, flags=re.DOTALL)
    slides = []
    for raw in re.split(r"\n---\n", md_text):
        raw = raw.strip()
        if not raw:
            continue
        lines = raw.split("\n")
        title, subtitle, body_lines, skip = "", "", [], False
        for i, line in enumerate(lines):
            if skip:
                skip = False
                continue
            if re.match(r"^# ", line):
                title = line[2:].strip()
                if i + 1 < len(lines) and re.match(r"^## ", lines[i + 1]):
                    subtitle = lines[i + 1][3:].strip()
                    skip = True
            elif re.match(r"^## ", line):
                title = line[3:].strip()
            else:
                body_lines.append(line)
        body = re.sub(r"<!--.*?-->", "", "\n".join(body_lines), flags=re.DOTALL).strip()
        if not (title or subtitle or body):
            continue   # üres-blokk gárda
        slides.append({"title": title, "subtitle": subtitle, "body": body,
                       "is_cover": bool(subtitle)})
    return slides


def parse_columns(body: str):
    """Kétoszlopos MARP div → (left_raw, right_img, right_cap, right_raw) | None.
    A bal/jobb szöveget NEM tisztítja — a szegmentálás a nyers szövegen dolgozik."""
    if '<div class="columns">' not in body:
        return None
    divs = re.findall(r"<div>(.*?)</div>", body, re.DOTALL)
    if len(divs) < 2:
        return None
    left_raw, right_raw = divs[0], divs[1]
    img_m = re.search(r"!\[[^\]]*\]\(([^)]+)\)", right_raw)
    cap_m = re.search(r'<span[^>]*class="cap"[^>]*>(.*?)</span>', right_raw, re.DOTALL)
    right_text = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", right_raw)
    right_text = re.sub(r'<span[^>]*class="cap"[^>]*>.*?</span>', "", right_text, flags=re.DOTALL)
    return (
        left_raw,
        img_m.group(1) if img_m else None,
        clean(cap_m.group(1)) if cap_m else "",
        right_text,
    )


# ── Routing-segédek ───────────────────────────────────────────────────────────

def hlevel_role(title: str) -> str:
    """Cím vezető száma → H1/H2/H3 ('2'→H1, '2.1'→H2, '2.1.1'→H3)."""
    num = nav.number_from_title(title)
    if not num:
        return "H2"
    return {0: "H1", 1: "H2", 2: "H3"}.get(num.count("."), "H2")


def resolve_current(title: str, nav_img: str | None, root) -> str | None:
    """Aktuális csomópont id-je: nav-képből (secN) vagy a cím számából."""
    if root is None:
        return None
    if nav_img:
        sec = nav.section_from_nav_image(nav_img)          # 'navigator'→None, secN→'N'
        node = nav.node_for_number(root, sec) if sec else None
        return node.id if node else None
    node = nav.node_for_number(root, nav.number_from_title(title))
    return node.id if node else None


def parse_gfm_table(table_lines):
    """GFM tábla-sorok → (headers, rows); a szeparátor-sort kihagyja."""
    headers, rows = [], []
    for line in table_lines:
        if re.match(r"^\s*\|[-: |]+\|\s*$", line):
            continue
        cells = [clean(c) for c in re.split(r"(?<!\\)\|", line) if c.strip() != ""]
        if not headers:
            headers = cells
        elif cells:
            rows.append(cells)
    return headers, rows


def add_pptx_table(slide, left, top, width, height, headers, rows):
    """Valódi PPTX-tábla beszúrása (DUE navy fejléc, Aptos)."""
    n_rows = len(rows) + 1
    n_cols = max(len(headers), max((len(r) for r in rows), default=0))
    if n_rows < 1 or n_cols < 1:
        return None
    row_h = min(int(height // n_rows), int(Emu(411480)))   # ≤ ~0.45"
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_h * n_rows).table
    NAVY, WHITE, DARK = RGBColor(0x0D, 0x1B, 0x3E), RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0x21, 0x21, 0x21)

    def setc(r, c, text, bold=False, fg=None, bg=None):
        cell = tbl.cell(r, c)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        if not run.text:
            run.text = text
        run.font.size = Pt(12)
        run.font.bold = bold
        run.font.name = "Aptos"
        if fg:
            run.font.color.rgb = fg
        if bg:
            cell.fill.solid(); cell.fill.fore_color.rgb = bg

    for c, h in enumerate(headers[:n_cols]):
        setc(0, c, h, bold=True, fg=WHITE, bg=NAVY)
    for r, row in enumerate(rows):
        for c in range(n_cols):
            setc(r + 1, c, row[c] if c < len(row) else "", fg=DARK)
    return tbl


# ── Képlet → PNG (conda base python: matplotlib mathtext) ─────────────────────

# ── Szöveg-kitöltés natív egyenletekkel (inline $...$ + block $$...$$) ─────────

def _clean_md(t: str) -> str:
    """Inline markdown-emfázis levétele (a `$...$` math érintetlen)."""
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", t)
    t = re.sub(r"(?<!\*)\*(?!\*)(.+?)\*", r"\1", t)
    t = re.sub(r"`(.+?)`", r"\1", t)
    return t


_BLOCK_EQ = re.compile(r"\$\$(.+?)\$\$", re.DOTALL)


def fill_text_frame(tf, raw_text: str):
    """A text-frame feltöltése bekezdésekként, NATÍV OMML-egyenletekkel:
    `$$...$$` (akár többsoros) → középre zárt block; `$...$` → szövegközi;
    egyébként sima (markdown-tisztított) szöveg. OMML hiányában a $-szöveg marad."""
    tf.clear()
    parts = _BLOCK_EQ.split(raw_text)        # páratlan indexek = block-képletek
    for i, part in enumerate(parts):
        if i % 2 == 1:                       # block egyenlet
            tex = " ".join(part.split())
            xml = _omml.block_paragraph_xml(tex)
            _omml.append_paragraph(tf, xml or _omml.plain_paragraph_xml(tex))
            continue
        for line in part.split("\n"):        # szöveg-blokk → sorok
            if not line.strip():
                continue
            lvl = _level(line)
            body = re.sub(r"^\s*>\s?", "", line)              # blockquote
            body = re.sub(r"^\s*[-*•]\s+", "", body)          # lista-jelölő
            body = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", body)  # kép-szintaxis
            body = re.sub(r"<[^>]+>", "", body)               # HTML tag
            body = _clean_md(body).strip()                    # bold/italic/code
            if not body:
                continue
            if _omml.has_inline(body):
                xml = _omml.inline_paragraph_xml(body, lvl)
                if xml:
                    _omml.append_paragraph(tf, xml); continue
            _omml.append_paragraph(tf, _omml.plain_paragraph_xml(body, lvl))
    _omml.drop_leading_empty(tf)


# ── Tartalom-kitöltés: szöveg+egyenlet a placeholderbe, tábla külön shape-ként ─

def _split_tables(raw_text: str):
    """Sorrendtartó bontás: ('text', str) és ('table', lines) blokkok."""
    blocks = []
    cur = []
    for line in raw_text.split("\n"):
        if re.match(r"^\s*\|", line):
            if cur:
                blocks.append(("text", "\n".join(cur))); cur = []
            if blocks and blocks[-1][0] == "table":
                blocks[-1][1].append(line)
            else:
                blocks.append(["table", [line]])
        else:
            cur.append(line)
    if cur:
        blocks.append(("text", "\n".join(cur)))
    return blocks


def render_content(slide, d, idx, raw_text, md_dir, formula_dir=None):
    """Az idx placeholderbe a próza + natív egyenletek (folyó szöveg); a markdown-
    táblák valódi PPTX-táblaként a régió alsó részébe kerülnek."""
    if idx not in d or not raw_text:
        return
    ph = d[idx]
    blocks = _split_tables(raw_text)
    tables = [b[1] for b in blocks if b[0] == "table"]

    if not tables:                                  # nincs tábla → minden folyó szöveg
        fill_text_frame(ph.text_frame, raw_text)
        return

    try:
        L, T, W, H = ph.left, ph.top, ph.width, ph.height
    except Exception:
        fill_text_frame(ph.text_frame, raw_text)
        return

    prose = "\n".join(b[1] for b in blocks if b[0] == "text" and b[1].strip()).strip()
    has_prose = bool(prose)
    text_h = int(H * 0.40) if has_prose else 0
    fill_text_frame(ph.text_frame, prose)
    if has_prose:
        try:                       # MIND a 4 dimenziót kiírjuk (különben W→0)
            ph.left, ph.top, ph.width, ph.height = L, T, W, text_h
        except Exception:
            pass
    gap = int(Emu(91440))
    vy = T + (text_h + gap if has_prose else 0)
    vh = int((H - (vy - T)) // max(1, len(tables)))
    for tbl in tables:
        headers, rows = parse_gfm_table(tbl)
        if headers:
            add_pptx_table(slide, L, vy, W, vh, headers, rows)
        vy += vh


# ── Dia-építők (variáns-tudatos) ──────────────────────────────────────────────

def _fill_title(d, sd, variant, root, current_id):
    """Default variánsban a többsoros breadcrumb CSAK alszakasz-diákon (N.M cím);
    a szakasz-zárók, áttekintő stb. megtartják a literál címüket."""
    if 0 not in d:
        return
    if variant == "default" and root is not None and current_id:
        node = nav.find_node(root, current_id)
        if node and node.num and "." in node.num:    # csak valódi alszakasz
            bc = nav.render_breadcrumb(root, current_id)
            if bc:
                set_title(d[0], bc)
                return
    set_title(d[0], sd["title"])


def _fill_sidebar(d, variant, root, current_id):
    if variant == "mindmap" and root is not None and 5 in d:
        set_tf(d[5], nav.render_toc(root, current_id))


def add_cover(prs, sd, variant):
    slide = prs.slides.add_slide(role_layout(prs, variant, "COVER"))
    d = slide_phs(slide)
    if 0 in d:
        set_title(d[0], sd["title"])
    if 1 in d and sd["subtitle"]:
        set_ph_text(d[1], sd["subtitle"])
    return "COVER"


def add_body(prs, sd, raw_body, role, variant, root, current_id, md_dir, formula_dir):
    slide = prs.slides.add_slide(role_layout(prs, variant, role))
    d = slide_phs(slide)
    _fill_title(d, sd, variant, root, current_id)
    render_content(slide, d, 1, raw_body, md_dir, formula_dir)
    _fill_sidebar(d, variant, root, current_id)
    return role


def add_kep_szoveg(prs, sd, left_raw, right_img, right_cap, md_dir, formula_dir,
                   variant, root, current_id):
    slide = prs.slides.add_slide(role_layout(prs, variant, "KEP"))
    d = slide_phs(slide)
    _fill_title(d, sd, variant, root, current_id)
    render_content(slide, d, 1, left_raw, md_dir, formula_dir)
    if 2 in d and right_img:
        insert_img_fit(slide, d[2], right_img, md_dir)
    if 3 in d and right_cap:
        set_ph_text(d[3], right_cap)
    _fill_sidebar(d, variant, root, current_id)
    return "KEP"


def add_section(prs, sd, desc_raw, variant, root, current_id, md_dir, formula_dir):
    """Szakasz-nyitó a Szakaszfejléc mintával: szám (idx1) + cím (idx0) + leírás (idx2)."""
    slide = prs.slides.add_slide(role_layout(prs, variant, "SECTION"))
    d = slide_phs(slide)
    title = sd["title"]
    num = nav.number_from_title(title) or ""
    title_text = re.sub(r"^\s*\d+\.\s*", "", title)
    if 0 in d:
        set_title(d[0], title_text)
    if 1 in d:
        set_ph_text(d[1], num)
    render_content(slide, d, 2, desc_raw, md_dir, formula_dir)
    _fill_sidebar(d, variant, root, current_id)
    return "SECTION"


def add_toc_overview(prs, sd, variant, root):
    """Default Áttekintés: a teljes hierarchikus TOC a body-ban (TOC layout)."""
    slide = prs.slides.add_slide(role_layout(prs, variant, "TOC"))
    d = slide_phs(slide)
    if 0 in d:
        set_title(d[0], sd["title"])
    if 1 in d and root is not None:
        set_tf(d[1], nav.render_toc(root, None))
    _fill_sidebar(d, variant, root, None)
    return "TOC"


# ── Fő builder ───────────────────────────────────────────────────────────────

def build_presentation(slides_data, potx_path, md_dir=None, *,
                       variant="default", mindmap_path=None):
    potx_path = Path(potx_path)
    if not potx_path.exists():
        sys.exit(f"Template nem található: {potx_path}")

    prs = load_potx(potx_path)
    md_dir = str(md_dir) if md_dir else None
    formula_dir = (str(Path(md_dir) / "_prezi_assets" / "_formulas")
                   if md_dir else "_formulas")
    root = (nav.parse_mindmap(mindmap_path)
            if mindmap_path and Path(mindmap_path).exists() else None)
    if root is None and variant == "mindmap":
        print("  ⚠️  Nincs mindmap.md — a mindmap variáns TOC nélkül készül.")
    print(f"Template: {potx_path.name}  ({len(prs.slide_layouts)} layout)  variáns={variant}")

    for i, sd in enumerate(slides_data):
        title, body = sd["title"], sd["body"]
        low = title.lower()
        cols = parse_columns(body) if '<div class="columns">' in body else None
        num = nav.number_from_title(title)
        is_top = bool(num) and "." not in num
        is_closer = ("összegzés" in low) or ("összefoglal" in low)

        if sd["is_cover"]:
            used = add_cover(prs, sd, variant)

        elif "hivatkozásjegyzék" in low or "irodalom" in low:
            used = add_body(prs, sd, body, "IROD", variant, root, None, md_dir, formula_dir)

        elif "áttekintés" in low:
            left_raw = cols[0] if cols else body
            if variant == "default":
                used = add_toc_overview(prs, sd, variant, root)
            else:
                used = add_body(prs, sd, left_raw, "TOC", variant, root, None, md_dir, formula_dir)

        elif is_top and not is_closer:
            # szakasz-NYITÓ → Szakaszfejléc minta
            cur = resolve_current(title, None, root)
            left_raw = cols[0] if cols else body
            used = add_section(prs, sd, left_raw, variant, root, cur, md_dir, formula_dir)

        elif cols:
            left_raw, right_img, right_cap, right_raw = cols
            if right_img and nav.is_nav_image(right_img):
                # navigációs kép → elhagyjuk; a navigáció idx5/idx0 felé megy
                cur = resolve_current(title, right_img, root)
                used = add_body(prs, sd, left_raw, hlevel_role(title), variant, root, cur, md_dir, formula_dir)
            elif right_img:
                cur = resolve_current(title, None, root)
                used = add_kep_szoveg(prs, sd, left_raw, right_img, right_cap,
                                      md_dir, formula_dir, variant, root, cur)
            else:
                cur = resolve_current(title, None, root)
                combined = (left_raw + "\n\n" + right_raw).strip()
                used = add_body(prs, sd, combined, hlevel_role(title), variant, root, cur, md_dir, formula_dir)

        else:
            cur = resolve_current(title, None, root)
            used = add_body(prs, sd, body, hlevel_role(title), variant, root, cur, md_dir, formula_dir)

        print(f"  [{i+1:2d}] {used:8s} ← {title[:42]}")

    return prs


# ── CLI ───────────────────────────────────────────────────────────────────────

def _detect_week(week_dir: Path) -> int:
    m = re.match(r"^(\d+)_het$", week_dir.name)
    if m:
        return int(m.group(1))
    raise ValueError(f"Nem tudja meghatározni a hét számát: {week_dir.name!r}")


def _out_for(base: Path, variant: str) -> Path:
    """default → base; mindmap → base_mindmap.pptx"""
    if variant == "default":
        return base
    return base.with_name(f"{base.stem}_{variant}{base.suffix}")


def run_one(md_path: Path, out_base: Path, template, mindmap_path, variant: str):
    md_text = md_path.read_text(encoding="utf-8")
    slides_data = parse_marp(md_text)
    potx = template or POTX_BY_VARIANT[variant]
    out_path = _out_for(out_base, variant)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    print(f"Feldolgozva: {len(slides_data)} dia  ←  {md_path}")
    prs = build_presentation(slides_data, potx, md_dir=md_path.parent,
                             variant=variant, mindmap_path=mindmap_path)
    prs.save(str(out_path))
    print(f"Mentve: {out_path}\n")


def main():
    parser = argparse.ArgumentParser(
        description="MARP MD → DUE PPTX (default / mindmap variáns)")
    parser.add_argument("input", nargs="?", default=None,
                        help="MARP .md fájl — elhagyható ha --week-dir megadva")
    parser.add_argument("--week-dir", default=None,
                        help="Pipeline hét-mappa (pl. test_outputs/<tantárgy>/N_het)")
    parser.add_argument("--variant", choices=["default", "mindmap", "both"],
                        default="default", help="Prezentáció-variáns (alap: default)")
    parser.add_argument("--template", default=None,
                        help="POTX template felülírás (alap: variáns szerint)")
    parser.add_argument("--mindmap", default=None,
                        help="mindmap.md útvonal (alap: <week>/3_mindmap/mindmap.md)")
    parser.add_argument("--output", default=None,
                        help="Kimeneti .pptx alap (elhagyható ha --week-dir megadva)")
    args = parser.parse_args()

    if args.week_dir:
        week_dir = Path(args.week_dir)
        if not week_dir.is_dir():
            sys.exit(f"Nem található: {week_dir}")
        n = _detect_week(week_dir)
        md_path = week_dir / "4_wip_outputs" / f"{n}_Prezentacio.md"
        out_base = week_dir / "5_clean_outputs" / f"{n}_Prezentacio.pptx"
        mindmap_path = args.mindmap or str(week_dir / "3_mindmap" / "mindmap.md")
        if args.output:
            out_base = Path(args.output)
    elif args.input:
        md_path = Path(args.input)
        out_base = Path(args.output) if args.output else md_path.with_suffix(".pptx")
        mindmap_path = args.mindmap
    else:
        parser.error("Kötelező: --week-dir VAGY positional input (md fájl)")

    if not md_path.exists():
        sys.exit(f"Nem található: {md_path}")

    template = Path(args.template) if args.template else None
    variants = ["default", "mindmap"] if args.variant == "both" else [args.variant]
    for v in variants:
        run_one(md_path, out_base, template, mindmap_path, v)


if __name__ == "__main__":
    main()
