"""Mini-env teszt 2: _omml.py valós kevert sorral + block képlettel."""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
import _omml

print("OMML available:", _omml.available())

prs = Presentation(); prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
tf = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(6)).text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "2.3. Biztonsági tartalék"
tf.paragraphs[0].runs[0].font.size = Pt(24)

lines = [
    "Mivel a surge vonal helye bizonytalan, biztonsági ráhagyást tartunk.",
    "ahol $P_o$ és $P_i$ a ki- és belépő össznyomás, $SA$ a surge-elkerülési pont.",
    "A B-paraméter $B = \\frac{U}{2a}\\sqrt{\\frac{V_p}{A_c L_c}}$ dönti el a típust.",
]
for ln in lines:
    xml = _omml.inline_paragraph_xml(ln)
    if xml:
        _omml.append_paragraph(tf, xml)
    else:
        p = tf.add_paragraph(); p.text = ln

# block
xml = _omml.block_paragraph_xml(r"SM = \frac{(P_o/P_i)_{Surge} - (P_o/P_i)_{SA}}{(P_o/P_i)_{SA}}")
_omml.append_paragraph(tf, xml)

out = Path(__file__).parent / "test2.pptx"
prs.save(str(out)); print("Mentve:", out)
